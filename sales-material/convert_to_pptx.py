#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""フラットサポート営業資料をPowerPointに変換"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# カラー定義
PRIMARY = RgbColor(26, 54, 93)
PRIMARY_LIGHT = RgbColor(44, 82, 130)
ACCENT = RgbColor(237, 137, 54)
ACCENT_LIGHT = RgbColor(246, 173, 85)
SKY_BLUE = RgbColor(66, 153, 225)
SUCCESS = RgbColor(72, 187, 120)
DANGER = RgbColor(229, 62, 62)
TEXT_DARK = RgbColor(26, 32, 44)
TEXT_MEDIUM = RgbColor(74, 85, 104)
TEXT_LIGHT = RgbColor(113, 128, 150)
WHITE = RgbColor(255, 255, 255)
BG_LIGHT = RgbColor(247, 250, 252)

def add_title_slide(prs, title, subtitle, stats=None):
    """表紙スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # バッジ
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(1.5), Inches(3), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].text = "初月500円キャンペーン中"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    tf.paragraphs[0].text = subtitle
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = RgbColor(200, 200, 200)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 統計
    if stats:
        y_pos = Inches(4.5)
        x_start = Inches(1.5)
        for i, (num, unit, label) in enumerate(stats):
            stat_box = slide.shapes.add_textbox(x_start + Inches(i * 2.5), y_pos, Inches(2), Inches(1.2))
            tf = stat_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = num
            run.font.size = Pt(48)
            run.font.color.rgb = ACCENT_LIGHT
            run.font.bold = True
            run2 = p.add_run()
            run2.text = unit
            run2.font.size = Pt(18)
            run2.font.color.rgb = ACCENT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            
            p2 = tf.add_paragraph()
            p2.text = label
            p2.font.size = Pt(11)
            p2.font.color.rgb = RgbColor(180, 180, 180)
            p2.alignment = PP_ALIGN.CENTER
    
    # 会社名
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    tf = company_box.text_frame
    tf.paragraphs[0].text = "株式会社UNISIA（ユニシア）"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = RgbColor(150, 150, 150)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_section_slide(prs, section_num, title, content_func):
    """セクションスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # セクション番号
    sec_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(2), Inches(0.3))
    tf = sec_box.text_frame
    tf.paragraphs[0].text = f"Section {section_num:02d}"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = ACCENT
    tf.paragraphs[0].font.bold = True
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].font.bold = True
    
    # 下線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # コンテンツを追加
    content_func(slide)
    
    return slide

def add_footer(slide, page_num, total_pages):
    """フッター"""
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "株式会社UNISIA"
    run1.font.size = Pt(10)
    run1.font.color.rgb = TEXT_LIGHT
    
    # ページ番号（右寄せ）
    page_box = slide.shapes.add_textbox(Inches(8.5), Inches(7.1), Inches(1), Inches(0.3))
    tf = page_box.text_frame
    tf.paragraphs[0].text = f"{page_num} / {total_pages}"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_pages = 15
    
    # スライド1: 表紙
    add_title_slide(prs, "フラットサポート", "月額制の航空券手配サービス", 
                    [("34", "万円", "最大節約額"), ("17", "カ国", "代表の渡航経験"), ("500", "円", "初月お試し価格")])
    
    # スライド2: こんなお悩み
    def content_problems(slide):
        problems = [
            ("😰", "航空券、高すぎない？"),
            ("😩", "自分で調べるのが面倒…"),
            ("😱", "家族4人分、いくらになるか怖い"),
            ("🤔", "エージェントに頼むと高いって聞くけど…"),
            ("😕", "本当に最安値か分からない"),
            ("😣", "相談できる人がいない"),
        ]
        y_start = Inches(1.6)
        for i, (emoji, text) in enumerate(problems):
            row = i // 2
            col = i % 2
            x = Inches(0.5) + Inches(col * 4.5)
            y = y_start + Inches(row * 1.0)
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.2), Inches(0.8))
            box.fill.solid()
            box.fill.fore_color.rgb = BG_LIGHT
            box.line.fill.background()
            
            tf = box.text_frame
            tf.paragraphs[0].text = f"{emoji}  {text}"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = TEXT_DARK
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            tf.word_wrap = True
        
        # 結論ボックス
        result = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.2), Inches(7), Inches(0.8))
        result.fill.solid()
        result.fill.fore_color.rgb = RgbColor(255, 243, 224)
        result.line.color.rgb = ACCENT
        tf = result.text_frame
        tf.paragraphs[0].text = "これらの悩み、すべて解決します。"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 2, total_pages)
    
    add_section_slide(prs, 1, "こんなお悩み、ありませんか？", content_problems)
    
    # スライド3: なぜ航空券は高いのか
    def content_why(slide):
        # 原因ボックス
        cause = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.8))
        cause.fill.solid()
        cause.fill.fore_color.rgb = RgbColor(254, 226, 226)
        cause.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(0.5))
        tf = title_box.text_frame
        tf.paragraphs[0].text = "答えはシンプル。「探し方を知らない」から。"
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.color.rgb = DANGER
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(8.6), Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = "多くの人は、大手比較サイトや留学エージェントの言い値で航空券を買います。\nでも実は、同じ目的地でも ルート・時期・航空会社の組み合わせ で価格は数万円〜数十万円変わります。"
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 知っていましたか
        know_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(9), Inches(1.5))
        know_box.fill.solid()
        know_box.fill.fore_color.rgb = RgbColor(237, 242, 247)
        know_box.line.fill.background()
        
        know_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.75), Inches(8.6), Inches(0.4))
        tf = know_title.text_frame
        tf.paragraphs[0].text = "知っていましたか？"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].font.bold = True
        
        points = [
            "• 留学エージェントは「最安値」ではなく「利益率の高い航空券」を提案している",
            "• 同じイギリス往復でも、取り方次第で 22万円の差 が出ることがある",
            "• 比較サイトには載っていない「裏ルート」が存在する"
        ]
        know_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(8.6), Inches(0.9))
        tf = know_text.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            if i == 0:
                tf.paragraphs[0].text = point
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
            else:
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(12)
                p.font.color.rgb = TEXT_MEDIUM
        
        # 解決ボックス
        solution = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(9), Inches(1.3))
        solution.fill.solid()
        solution.fill.fore_color.rgb = SKY_BLUE
        solution.line.fill.background()
        
        sol_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(8.6), Inches(0.4))
        tf = sol_title.text_frame
        tf.paragraphs[0].text = "だから、私たちが代わりに探します。"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        sol_desc = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(8.6), Inches(0.5))
        tf = sol_desc.text_frame
        tf.paragraphs[0].text = "17カ国渡航の経験を持つ代表が、あなたの航空券を「最安値」で見つけます。"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RgbColor(220, 220, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 3, total_pages)
    
    add_section_slide(prs, 2, "なぜ、航空券は高いのか？", content_why)
    
    # スライド4: フラットサポートとは
    def content_what(slide):
        # メインボックス
        main_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.6), Inches(7), Inches(1.2))
        main_box.fill.solid()
        main_box.fill.fore_color.rgb = RgbColor(255, 243, 224)
        main_box.line.color.rgb = ACCENT
        
        main_text = slide.shapes.add_textbox(Inches(1.7), Inches(1.8), Inches(6.6), Inches(0.9))
        tf = main_text.text_frame
        tf.paragraphs[0].text = "業界初のサブスクリプション型"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "航空券手配サービス"
        p.font.size = Pt(26)
        p.font.color.rgb = ACCENT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # 説明
        desc = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
        tf = desc.text_frame
        tf.paragraphs[0].text = "エージェントのようにパッケージを売りつけません。\nシンプルに「航空券を最安値で取る」ことだけに集中します。"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # STEP
        steps = [
            ("STEP 1", "LINEで「○月に○○に行きたい」と送る"),
            ("STEP 2", "私たちが最安値のルートを探す"),
            ("STEP 3", "複数の選択肢をご提案"),
            ("STEP 4", "お好みのプランで手配完了"),
        ]
        y_start = Inches(3.8)
        for i, (step, desc_text) in enumerate(steps):
            step_box = slide.shapes.add_textbox(Inches(1), y_start + Inches(i * 0.5), Inches(2), Inches(0.4))
            tf = step_box.text_frame
            tf.paragraphs[0].text = step
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = PRIMARY
            tf.paragraphs[0].font.bold = True
            
            desc_box = slide.shapes.add_textbox(Inches(3), y_start + Inches(i * 0.5), Inches(6), Inches(0.4))
            tf = desc_box.text_frame
            tf.paragraphs[0].text = desc_text
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        
        # 結論
        result = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.9), Inches(7), Inches(0.7))
        result.fill.solid()
        result.fill.fore_color.rgb = RgbColor(237, 242, 247)
        result.line.fill.background()
        
        result_text = slide.shapes.add_textbox(Inches(1.7), Inches(6.05), Inches(6.6), Inches(0.4))
        tf = result_text.text_frame
        tf.paragraphs[0].text = "LINEで相談するだけ。あとは全部おまかせ。"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 4, total_pages)
    
    add_section_slide(prs, 3, "フラットサポートとは", content_what)
    
    # スライド5: かんたん操作ガイド
    def content_guide(slide):
        # 説明
        desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
        tf = desc.text_frame
        tf.paragraphs[0].text = "LINEのリッチメニューから「格安航空券サポート」をタップするだけ。\nその日の最安値を、リアルタイムでご提案します。"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 4ステップ（画像付き）
        steps = [
            ("STEP 1", "リッチメニューをタップ", "「格安航空券サポート」を選択", "images/step2-richmenu.png"),
            ("STEP 2", "テンプレートに入力", "行きたい地域・時期・人数など", "images/step1-template.png"),
            ("STEP 3", "最安値フライトを表示", "リアルタイムの格安航空券", "images/step3-flight.png"),
            ("STEP 4", "ホテルも検索可能", "オプションで条件指定OK", "images/step4-hotel.png"),
        ]
        
        script_dir = os.path.dirname(os.path.abspath(__file__))
        
        for i, (step, title, desc_text, img_path) in enumerate(steps):
            x = Inches(0.4) + Inches(i * 2.4)
            y = Inches(2.2)
            
            # STEPラベル
            if i < 3:
                label_color = SKY_BLUE
            else:
                label_color = ACCENT
            
            label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y, Inches(1.5), Inches(0.35))
            label.fill.solid()
            label.fill.fore_color.rgb = label_color
            label.line.fill.background()
            tf = label.text_frame
            tf.paragraphs[0].text = step
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 画像
            full_img_path = os.path.join(script_dir, img_path)
            if os.path.exists(full_img_path):
                try:
                    pic = slide.shapes.add_picture(full_img_path, x, y + Inches(0.45), width=Inches(2.2))
                    # アスペクト比を維持しながら高さを調整
                    if pic.height > Inches(2.5):
                        ratio = Inches(2.5) / pic.height
                        pic.height = Inches(2.5)
                        pic.width = int(pic.width * ratio)
                except:
                    # 画像が読み込めない場合はプレースホルダー
                    placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.45), Inches(2.1), Inches(2.5))
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = BG_LIGHT
                    placeholder.line.color.rgb = TEXT_LIGHT
            else:
                placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.45), Inches(2.1), Inches(2.5))
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = BG_LIGHT
                placeholder.line.color.rgb = TEXT_LIGHT
            
            # タイトル
            title_box = slide.shapes.add_textbox(x, y + Inches(3.1), Inches(2.2), Inches(0.35))
            tf = title_box.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = PRIMARY
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 説明
            desc_box = slide.shapes.add_textbox(x, y + Inches(3.4), Inches(2.2), Inches(0.3))
            tf = desc_box.text_frame
            tf.paragraphs[0].text = desc_text
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 下部のポイント
        solution = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(9), Inches(0.8))
        solution.fill.solid()
        solution.fill.fore_color.rgb = SKY_BLUE
        solution.line.fill.background()
        
        sol_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(8.6), Inches(0.6))
        tf = sol_text.text_frame
        tf.paragraphs[0].text = "「格安サイトで最安値のつもり」が実は最安値じゃなかった…"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "フラットサポートなら、その場のリアルタイムレートで本当の最安値をご提案します。"
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(220, 220, 255)
        p.alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 5, total_pages)
    
    add_section_slide(prs, 4, "かんたん操作ガイド", content_guide)
    
    # スライド6: 選ばれる3つの理由
    def content_reasons(slide):
        reasons = [
            ("💰", "POINT 01", "最大34万円安くなる", "大手エージェントより圧倒的に安く航空券を手配。\n世界一周16万円の実績あり。"),
            ("👨‍👩‍👧‍👦", "POINT 02", "家族全員分サポート", "月額8,800円で、ご家族全員の航空券を手配。\n追加料金は一切かかりません。"),
            ("💬", "POINT 03", "LINEでいつでも相談", "「この時期はいくら？」「どのルートが安い？」\n気軽にLINEで相談できます。"),
        ]
        
        for i, (emoji, point, title, desc) in enumerate(reasons):
            x = Inches(0.5) + Inches(i * 3.1)
            y = Inches(1.6)
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(3.5))
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = BG_LIGHT
            
            # 絵文字
            emoji_box = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2.9), Inches(0.6))
            tf = emoji_box.text_frame
            tf.paragraphs[0].text = emoji
            tf.paragraphs[0].font.size = Pt(36)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # POINT
            point_box = slide.shapes.add_textbox(x, y + Inches(0.85), Inches(2.9), Inches(0.3))
            tf = point_box.text_frame
            tf.paragraphs[0].text = point
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = ACCENT
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # タイトル
            title_box = slide.shapes.add_textbox(x, y + Inches(1.15), Inches(2.9), Inches(0.4))
            tf = title_box.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = PRIMARY
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 説明
            desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.65), Inches(2.7), Inches(1.5))
            tf = desc_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = desc
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 500円ボックス
        trial = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.4), Inches(7), Inches(1.0))
        trial.fill.solid()
        trial.fill.fore_color.rgb = RgbColor(255, 243, 224)
        trial.line.color.rgb = ACCENT
        
        trial_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.55), Inches(6.6), Inches(0.7))
        tf = trial_text.text_frame
        tf.paragraphs[0].text = "さらに、初月500円でお試しいただけます。"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "合わなければ即解約OK。解約手数料なし。"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 6, total_pages)
    
    add_section_slide(prs, 5, "選ばれる3つの理由", content_reasons)
    
    # スライド7: 実際にお得になった事例
    def content_cases(slide):
        cases = [
            ("CASE 01", "🌍 世界一周航空券", "50万円", "16万円", "34万円お得！"),
            ("CASE 02", "🇬🇧 イギリス往復", "33万円", "11万円", "22万円お得！"),
            ("CASE 03", "🇦🇺 オーストラリア往復", "18万円", "12万円", "6万円お得！"),
        ]
        
        for i, (case, dest, before, after, saving) in enumerate(cases):
            x = Inches(0.5) + Inches(i * 3.1)
            y = Inches(1.6)
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.9), Inches(2.8))
            box.fill.solid()
            box.fill.fore_color.rgb = WHITE
            box.line.color.rgb = BG_LIGHT
            
            # CASE
            case_box = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(2.9), Inches(0.3))
            tf = case_box.text_frame
            tf.paragraphs[0].text = case
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            case_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.8), y + Inches(0.1), Inches(1.3), Inches(0.3))
            case_bg.fill.solid()
            case_bg.fill.fore_color.rgb = PRIMARY
            case_bg.line.fill.background()
            case_bg.text_frame.paragraphs[0].text = case
            case_bg.text_frame.paragraphs[0].font.size = Pt(9)
            case_bg.text_frame.paragraphs[0].font.color.rgb = WHITE
            case_bg.text_frame.paragraphs[0].font.bold = True
            case_bg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 行き先
            dest_box = slide.shapes.add_textbox(x, y + Inches(0.5), Inches(2.9), Inches(0.4))
            tf = dest_box.text_frame
            tf.paragraphs[0].text = dest
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = TEXT_DARK
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Before
            before_box = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(2.9), Inches(0.5))
            tf = before_box.text_frame
            tf.paragraphs[0].text = "大手エージェント"
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.add_paragraph()
            p.text = before
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            
            # 矢印
            arrow_box = slide.shapes.add_textbox(x, y + Inches(1.55), Inches(2.9), Inches(0.3))
            tf = arrow_box.text_frame
            tf.paragraphs[0].text = "↓"
            tf.paragraphs[0].font.size = Pt(16)
            tf.paragraphs[0].font.color.rgb = ACCENT
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # After
            after_box = slide.shapes.add_textbox(x, y + Inches(1.8), Inches(2.9), Inches(0.5))
            tf = after_box.text_frame
            tf.paragraphs[0].text = "フラットサポート"
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = ACCENT
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = tf.add_paragraph()
            p.text = after
            p.font.size = Pt(18)
            p.font.color.rgb = PRIMARY
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # 節約額
            saving_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.4), y + Inches(2.4), Inches(2.1), Inches(0.35))
            saving_bg.fill.solid()
            saving_bg.fill.fore_color.rgb = ACCENT
            saving_bg.line.fill.background()
            saving_bg.text_frame.paragraphs[0].text = saving
            saving_bg.text_frame.paragraphs[0].font.size = Pt(12)
            saving_bg.text_frame.paragraphs[0].font.color.rgb = WHITE
            saving_bg.text_frame.paragraphs[0].font.bold = True
            saving_bg.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 6ヶ月留学の場合
        compare_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.7), Inches(9), Inches(1.8))
        compare_box.fill.solid()
        compare_box.fill.fore_color.rgb = RgbColor(237, 242, 247)
        compare_box.line.fill.background()
        
        compare_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(8.6), Inches(0.35))
        tf = compare_title.text_frame
        tf.paragraphs[0].text = "6ヶ月留学の場合"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].font.bold = True
        
        # 比較テーブル（テキストベース）
        table_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.3), Inches(8.6), Inches(1.0))
        tf = table_text.text_frame
        tf.paragraphs[0].text = "大手留学エージェント: 約190万円〜　→　UNISIA: 65万円〜"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "約125万円お得"
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 7, total_pages)
    
    add_section_slide(prs, 6, "実際にお得になった事例", content_cases)
    
    # スライド8: 料金プラン
    def content_pricing(slide):
        # メインプラン
        main_plan = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(5.5), Inches(4.8))
        main_plan.fill.solid()
        main_plan.fill.fore_color.rgb = WHITE
        main_plan.line.color.rgb = ACCENT
        main_plan.line.width = Pt(3)
        
        # 人気No.1バッジ
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), Inches(1.45), Inches(1.9), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = "人気No.1"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # プラン名
        plan_name = slide.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(5.1), Inches(0.4))
        tf = plan_name.text_frame
        tf.paragraphs[0].text = "フラットサポート"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        plan_desc = slide.shapes.add_textbox(Inches(0.7), Inches(2.35), Inches(5.1), Inches(0.3))
        tf = plan_desc.text_frame
        tf.paragraphs[0].text = "航空券手配 + フルサポート"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 価格
        price_label = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(5.1), Inches(0.25))
        tf = price_label.text_frame
        tf.paragraphs[0].text = "初月"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        price = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(5.1), Inches(0.7))
        tf = price.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "500"
        run.font.size = Pt(48)
        run.font.color.rgb = ACCENT
        run.font.bold = True
        run2 = p.add_run()
        run2.text = "円"
        run2.font.size = Pt(18)
        run2.font.color.rgb = TEXT_MEDIUM
        p.alignment = PP_ALIGN.CENTER
        
        price_after = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(5.1), Inches(0.25))
        tf = price_after.text_frame
        tf.paragraphs[0].text = "翌月以降 8,800円/月"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 特徴リスト
        features = [
            "✓ 航空券の最安値検索＆手配",
            "✓ 家族全員分サポート（追加料金なし）",
            "✓ 24時間LINEサポート",
            "✓ 現地トラブル対応",
            "✓ ビザ申請サポート",
            "✓ 海外保険のご案内",
        ]
        y_start = Inches(4.1)
        for i, feature in enumerate(features):
            f_box = slide.shapes.add_textbox(Inches(1.0), y_start + Inches(i * 0.32), Inches(4.5), Inches(0.3))
            tf = f_box.text_frame
            tf.paragraphs[0].text = feature
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        
        # プレミアムプラン
        premium = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(1.6), Inches(3.3), Inches(4.8))
        premium.fill.solid()
        premium.fill.fore_color.rgb = WHITE
        premium.line.color.rgb = BG_LIGHT
        
        prem_name = slide.shapes.add_textbox(Inches(6.4), Inches(1.85), Inches(2.9), Inches(0.4))
        tf = prem_name.text_frame
        tf.paragraphs[0].text = "プレミアム"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        prem_desc = slide.shapes.add_textbox(Inches(6.4), Inches(2.25), Inches(2.9), Inches(0.3))
        tf = prem_desc.text_frame
        tf.paragraphs[0].text = "留学・ワーホリをフルサポート"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        prem_price = slide.shapes.add_textbox(Inches(6.4), Inches(2.7), Inches(2.9), Inches(0.5))
        tf = prem_price.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "19,800"
        run.font.size = Pt(28)
        run.font.color.rgb = PRIMARY
        run.font.bold = True
        run2 = p.add_run()
        run2.text = "円/月"
        run2.font.size = Pt(12)
        run2.font.color.rgb = TEXT_MEDIUM
        p.alignment = PP_ALIGN.CENTER
        
        prem_features = [
            "✓ フラットサポートの内容すべて",
            "✓ ビザ申請代行",
            "✓ オンライン英会話（月4回）",
            "✓ 帰国後の就職サポート",
        ]
        y_start = Inches(3.4)
        for i, feature in enumerate(prem_features):
            f_box = slide.shapes.add_textbox(Inches(6.4), y_start + Inches(i * 0.32), Inches(2.9), Inches(0.3))
            tf = f_box.text_frame
            tf.paragraphs[0].text = feature
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        
        add_footer(slide, 8, total_pages)
    
    add_section_slide(prs, 7, "料金プラン", content_pricing)
    
    # スライド9-15: 残りのスライドを追加
    # スライド9: 留学エージェントとの違い
    def content_compare(slide):
        items = [
            ("対象", "留学する人のみ", "✓ 旅行でも留学でもOK"),
            ("料金体系", "パッケージ（高額）", "✓ シンプルな月額制"),
            ("航空券", "割高なことが多い", "✓ 最安値を徹底追求"),
            ("家族利用", "追加料金あり", "✓ 追加料金なし"),
            ("相談方法", "予約制の面談", "✓ LINEでいつでも"),
            ("解約", "違約金あり", "✓ いつでも無料で解約"),
        ]
        
        y_start = Inches(1.7)
        # ヘッダー
        headers = ["比較項目", "留学エージェント", "フラットサポート"]
        header_colors = [TEXT_MEDIUM, TEXT_LIGHT, ACCENT]
        widths = [Inches(2.2), Inches(3.2), Inches(3.6)]
        x_pos = Inches(0.5)
        for i, (header, color, width) in enumerate(zip(headers, header_colors, widths)):
            h_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y_start, width, Inches(0.4))
            h_box.fill.solid()
            h_box.fill.fore_color.rgb = color if i < 2 else ACCENT
            h_box.line.fill.background()
            h_box.text_frame.paragraphs[0].text = header
            h_box.text_frame.paragraphs[0].font.size = Pt(11)
            h_box.text_frame.paragraphs[0].font.color.rgb = WHITE
            h_box.text_frame.paragraphs[0].font.bold = True
            h_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            x_pos += width
        
        # 行
        for row_idx, (item, agent, flat) in enumerate(items):
            y = y_start + Inches(0.45) + Inches(row_idx * 0.45)
            bg_color = WHITE if row_idx % 2 == 0 else BG_LIGHT
            
            x_pos = Inches(0.5)
            for col_idx, (text, width) in enumerate(zip([item, agent, flat], widths)):
                cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, y, width, Inches(0.42))
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                cell.line.fill.background()
                
                text_color = TEXT_DARK if col_idx == 0 else (TEXT_LIGHT if col_idx == 1 else SUCCESS)
                font_bold = col_idx == 0 or col_idx == 2
                
                cell.text_frame.paragraphs[0].text = text
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].font.color.rgb = text_color
                cell.text_frame.paragraphs[0].font.bold = font_bold
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT
                x_pos += width
        
        # 結論
        result = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(7), Inches(0.9))
        result.fill.solid()
        result.fill.fore_color.rgb = RgbColor(255, 243, 224)
        result.line.color.rgb = ACCENT
        
        result_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.15), Inches(6.6), Inches(0.6))
        tf = result_text.text_frame
        tf.paragraphs[0].text = "「航空券を安く取りたい」だけなら、"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "フラットサポートで十分です。"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 9, total_pages)
    
    add_section_slide(prs, 8, "留学エージェントとの違い", content_compare)
    
    # スライド10: お客様の声
    def content_voice(slide):
        voices = [
            ("Y", "male", "Y.Tさん（20代・男性）", "世界一周で利用", "世界一周の航空券が合計16万円で取れました！大手だと50万円近くかかると言われていたので、正直信じられなかったです。", True),
            ("M", "female", "M.Sさん（20代・女性）", "ワーホリ準備で利用", "オーストラリアへのワーホリ。航空券だけで6万円浮いて、その分を現地での生活費に回せました。", False),
            ("K", "male", "K.Nさん（30代・男性）", "家族旅行で利用", "家族4人でハワイに行きました。自分で調べた価格より8万円も安く取れました。LINEで相談するだけで全部やってくれるのが助かりました！", False),
            ("R", "female", "R.Kさん（20代・女性）", "LINEサポートが決め手", "何よりLINEでのサポートが手厚いのが良かったです！質問にもすぐ返信がきて、初めての海外でも安心して準備できました。", False),
        ]
        
        for i, (initial, gender, name, tag, text, highlight) in enumerate(voices):
            row = i // 2
            col = i % 2
            x = Inches(0.5) + Inches(col * 4.7)
            y = Inches(1.6) + Inches(row * 2.5)
            
            bg_color = PRIMARY if highlight else BG_LIGHT
            text_color = WHITE if highlight else TEXT_MEDIUM
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.4), Inches(2.3))
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
            box.line.fill.background()
            
            # アバター
            avatar_color = SKY_BLUE if gender == "male" else RgbColor(237, 100, 166)
            avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = avatar_color
            avatar.line.fill.background()
            avatar.text_frame.paragraphs[0].text = initial
            avatar.text_frame.paragraphs[0].font.size = Pt(16)
            avatar.text_frame.paragraphs[0].font.color.rgb = WHITE
            avatar.text_frame.paragraphs[0].font.bold = True
            avatar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 名前
            name_box = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.15), Inches(3.5), Inches(0.3))
            tf = name_box.text_frame
            tf.paragraphs[0].text = name
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = text_color
            tf.paragraphs[0].font.bold = True
            
            tag_box = slide.shapes.add_textbox(x + Inches(0.75), y + Inches(0.4), Inches(3.5), Inches(0.25))
            tf = tag_box.text_frame
            tf.paragraphs[0].text = tag
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = RgbColor(180, 180, 180) if highlight else TEXT_LIGHT
            
            # 星
            star_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.75), Inches(4.1), Inches(0.3))
            tf = star_box.text_frame
            tf.paragraphs[0].text = "★★★★★"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = ACCENT_LIGHT if highlight else ACCENT
            
            # テキスト
            text_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.05), Inches(4.1), Inches(1.1))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = text
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = text_color
        
        add_footer(slide, 10, total_pages)
    
    add_section_slide(prs, 9, "ご利用者様の声", content_voice)
    
    # スライド11: 代表紹介
    def content_about(slide):
        # 写真プレースホルダー
        photo_placeholder = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(2.2), Inches(2.2))
        photo_placeholder.fill.solid()
        photo_placeholder.fill.fore_color.rgb = SKY_BLUE
        photo_placeholder.line.fill.background()
        
        photo_text = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(2.2), Inches(0.5))
        tf = photo_text.text_frame
        tf.paragraphs[0].text = "PHOTO"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 17カ国バッジ
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(4.4), Inches(1.6), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = "🌍 17カ国渡航"
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # プロフィールテキスト
        name_box = slide.shapes.add_textbox(Inches(4), Inches(1.8), Inches(5.5), Inches(0.5))
        tf = name_box.text_frame
        tf.paragraphs[0].text = "トライ"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.bold = True
        
        title_box = slide.shapes.add_textbox(Inches(4), Inches(2.3), Inches(5.5), Inches(0.3))
        tf = title_box.text_frame
        tf.paragraphs[0].text = "株式会社UNISIA 代表取締役"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        
        story_box = slide.shapes.add_textbox(Inches(4), Inches(2.8), Inches(5.5), Inches(3.5))
        tf = story_box.text_frame
        tf.word_wrap = True
        
        story_text = """貯金5万円・片道切符で海外に飛び出し、17カ国以上を回りました。

一方で、留学エージェントに200万円損した経験も。高額なパッケージを売りつけられ、後から「もっと安くできた」と知りました。

あの時の自分と同じ思いをしてほしくない。だから、フラットサポートを作りました。

「航空券の取り方を変えるだけで、人生の選択肢が広がる。」
あなたの海外挑戦を、全力でサポートします。"""
        
        tf.paragraphs[0].text = story_text
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        tf.paragraphs[0].line_spacing = 1.5
        
        add_footer(slide, 11, total_pages)
    
    add_section_slide(prs, 10, "代表紹介", content_about)
    
    # スライド12: よくある質問
    def content_faq(slide):
        faqs = [
            ("本当に安くなりますか？", "はい。大手エージェントと比較して、最大34万円お得になった実績があります。航空券は時期やルートによって価格が変動しますが、常に最安値を徹底的にお探しします。"),
            ("家族は何人まで使えますか？", "1契約につき最大9名様までご利用いただけます。ご家族全員分を追加料金なしで手配しますので、家族旅行にもおすすめです。"),
            ("解約はいつでもできますか？", "はい。いつでも解約可能です。解約手数料もかかりません。合わなければ、1ヶ月で辞めていただいて構いません。"),
            ("航空券以外も相談できますか？", "はい。ビザ申請や海外保険のご案内、渡航中の現地トラブル対応、海外の情勢・治安情報の配信など、LINEでいつでも何度でもご相談いただけます。"),
        ]
        
        y_start = Inches(1.6)
        for i, (q, a) in enumerate(faqs):
            y = y_start + Inches(i * 1.35)
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(9), Inches(1.25))
            box.fill.solid()
            box.fill.fore_color.rgb = BG_LIGHT
            box.line.fill.background()
            
            # Q
            q_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), y + Inches(0.12), Inches(0.35), Inches(0.35))
            q_badge.fill.solid()
            q_badge.fill.fore_color.rgb = SKY_BLUE
            q_badge.line.fill.background()
            q_badge.text_frame.paragraphs[0].text = "Q"
            q_badge.text_frame.paragraphs[0].font.size = Pt(12)
            q_badge.text_frame.paragraphs[0].font.color.rgb = WHITE
            q_badge.text_frame.paragraphs[0].font.bold = True
            q_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            q_text = slide.shapes.add_textbox(Inches(1.1), y + Inches(0.12), Inches(8.2), Inches(0.35))
            tf = q_text.text_frame
            tf.paragraphs[0].text = q
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = TEXT_DARK
            tf.paragraphs[0].font.bold = True
            
            # A
            a_badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), y + Inches(0.55), Inches(0.35), Inches(0.35))
            a_badge.fill.solid()
            a_badge.fill.fore_color.rgb = ACCENT
            a_badge.line.fill.background()
            a_badge.text_frame.paragraphs[0].text = "A"
            a_badge.text_frame.paragraphs[0].font.size = Pt(12)
            a_badge.text_frame.paragraphs[0].font.color.rgb = WHITE
            a_badge.text_frame.paragraphs[0].font.bold = True
            a_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            a_text = slide.shapes.add_textbox(Inches(1.1), y + Inches(0.55), Inches(8.2), Inches(0.6))
            tf = a_text.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = a
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = TEXT_MEDIUM
        
        add_footer(slide, 12, total_pages)
    
    add_section_slide(prs, 11, "よくある質問", content_faq)
    
    # スライド13: サービス一覧
    def content_services(slide):
        services = [
            ("✈️", "航空券の手配", "最安値を徹底追求"),
            ("🚗", "空港送迎手配", "到着後も安心"),
            ("📄", "ビザ発行代行", "面倒な手続きを代行"),
            ("🛡️", "海外保険案内", "最適なプランをご提案"),
            ("💬", "24時間LINEサポート", "いつでも相談OK"),
            ("🆘", "現地トラブル対応", "渡航中も安心"),
            ("🗣️", "現地日本語サポート", "言葉の壁も安心"),
            ("🎓", "オンライン英会話", "プレミアムプランのみ"),
            ("💼", "帰国後就職サポート", "キャリアもサポート"),
            ("📷", "海外現地撮影", "思い出を形に"),
        ]
        
        for i, (emoji, title, desc) in enumerate(services):
            row = i // 2
            col = i % 2
            x = Inches(0.5) + Inches(col * 4.7)
            y = Inches(1.6) + Inches(row * 0.95)
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.4), Inches(0.85))
            box.fill.solid()
            box.fill.fore_color.rgb = BG_LIGHT
            box.line.fill.background()
            
            emoji_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), Inches(0.5), Inches(0.5))
            tf = emoji_box.text_frame
            tf.paragraphs[0].text = emoji
            tf.paragraphs[0].font.size = Pt(20)
            
            title_box = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.15), Inches(3.5), Inches(0.35))
            tf = title_box.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = PRIMARY
            tf.paragraphs[0].font.bold = True
            
            desc_box = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.45), Inches(3.5), Inches(0.3))
            tf = desc_box.text_frame
            tf.paragraphs[0].text = desc
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = TEXT_LIGHT
        
        add_footer(slide, 13, total_pages)
    
    add_section_slide(prs, 12, "提供サポート一覧", content_services)
    
    # スライド14: まとめ・CTA
    def content_cta(slide):
        # CTAボックス
        cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(4.0))
        cta.fill.solid()
        cta.fill.fore_color.rgb = PRIMARY
        cta.line.fill.background()
        
        cta_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(0.5))
        tf = cta_title.text_frame
        tf.paragraphs[0].text = "まずは初月500円でお試しください"
        tf.paragraphs[0].font.size = Pt(26)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        cta_sub = slide.shapes.add_textbox(Inches(0.7), Inches(2.35), Inches(8.6), Inches(0.35))
        tf = cta_sub.text_frame
        tf.paragraphs[0].text = "航空券1回の手配で、数万円〜数十万円の差が出ます。"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = RgbColor(200, 200, 200)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 価格
        price_label = slide.shapes.add_textbox(Inches(0.7), Inches(2.9), Inches(8.6), Inches(0.3))
        tf = price_label.text_frame
        tf.paragraphs[0].text = "初月"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RgbColor(150, 150, 150)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        price = slide.shapes.add_textbox(Inches(0.7), Inches(3.15), Inches(8.6), Inches(0.9))
        tf = price.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "500"
        run.font.size = Pt(64)
        run.font.color.rgb = ACCENT_LIGHT
        run.font.bold = True
        run2 = p.add_run()
        run2.text = "円"
        run2.font.size = Pt(22)
        run2.font.color.rgb = ACCENT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
        # 特徴
        features = ["家族全員分OK", "LINEでいつでも相談", "いつでも解約可能"]
        features_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(8.6), Inches(0.35))
        tf = features_text.text_frame
        tf.paragraphs[0].text = "　✓ " + "　　✓ ".join(features)
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 注釈
        note = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8.6), Inches(0.3))
        tf = note.text_frame
        tf.paragraphs[0].text = "※ 翌月以降 8,800円/月 ・ 解約手数料なし"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = RgbColor(150, 150, 150)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 下部ボックス
        bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.7), Inches(7), Inches(0.9))
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = RgbColor(255, 243, 224)
        bottom.line.color.rgb = ACCENT
        
        bottom_text = slide.shapes.add_textbox(Inches(1.7), Inches(5.85), Inches(6.6), Inches(0.6))
        tf = bottom_text.text_frame
        tf.paragraphs[0].text = "リスクはゼロ。合わなければ、500円以上の損はしません。"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = TEXT_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.add_paragraph()
        p.text = "でも、もし合えば、次の旅行から数万円〜数十万円浮くかもしれません。"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MEDIUM
        p.alignment = PP_ALIGN.CENTER
        
        add_footer(slide, 14, total_pages)
    
    add_section_slide(prs, 13, "まとめ", content_cta)
    
    # スライド15: お問い合わせ（最終ページ）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "お問い合わせ"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 連絡先
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3.2), Inches(6), Inches(2.5))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = RgbColor(44, 82, 130)
    contact_box.line.fill.background()
    
    contacts = [
        ("会社名", "株式会社UNISIA（ユニシア）"),
        ("代表取締役", "井上 智羅（いのうえ とらい）"),
        ("所在地", "福岡県福岡市博多区博多駅前1丁目23番2号"),
        ("サービスURL", "ltdunisia.memberpay.jp"),
        ("公式LINE", "@unisia"),
    ]
    
    y_start = Inches(3.4)
    for i, (label, value) in enumerate(contacts):
        row = slide.shapes.add_textbox(Inches(2.3), y_start + Inches(i * 0.42), Inches(5.4), Inches(0.4))
        tf = row.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{label}　　"
        run1.font.size = Pt(11)
        run1.font.color.rgb = RgbColor(150, 150, 180)
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(11)
        run2.font.color.rgb = WHITE
    
    # タグライン
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = tagline.text_frame
    tf.paragraphs[0].text = '"すべての人に、気軽な海外挑戦を"'
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = RgbColor(180, 180, 200)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "営業資料_フラットサポート.pptx")
    prs.save(output_path)
    print(f"PowerPointファイルを作成しました: {output_path}")
