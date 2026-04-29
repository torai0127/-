#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
フラットサポート営業資料 PowerPoint作成スクリプト v2
- レイアウト修正
- コミュニティ内容追加
- サポート内容詳細化
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラー定義
PRIMARY = RGBColor(26, 54, 93)
ACCENT = RGBColor(237, 137, 54)
ACCENT_LIGHT = RGBColor(246, 173, 85)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(26, 32, 44)
TEXT_MEDIUM = RGBColor(74, 85, 104)
TEXT_LIGHT = RGBColor(113, 128, 150)
SKY_BLUE = RGBColor(66, 153, 225)
SUCCESS = RGBColor(72, 187, 120)
DANGER = RGBColor(229, 62, 62)
BG_LIGHT = RGBColor(247, 250, 252)
BG_CREAM = RGBColor(255, 250, 240)
PINK_LIGHT = RGBColor(254, 215, 215)

def add_text_frame(shape, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    """シェイプ内にテキストを設定"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans JP"
    try:
        tf.vertical_anchor = valign
    except:
        pass
    return tf

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans JP"
    p.alignment = align
    return box

def add_rounded_box(slide, left, top, width, height, fill_color, line_color=None):
    """角丸ボックスを追加"""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if line_color:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box

def add_rect(slide, left, top, width, height, fill_color):
    """四角形を追加"""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    rect.line.fill.background()
    return rect

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ===== スライド1: 表紙 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.3, 10, 0.4, "初月500円キャンペーン中", 14, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.8, 10, 0.3, "LTD. UNISIA", 12, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.6, 10, 0.8, "フラットサポート", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.5, 10, 0.5, "月額制の航空券手配 & 海外挑戦サポート", 18, False, WHITE, PP_ALIGN.CENTER)

# 統計ボックス
stats = [("34", "万円", "最大節約額"), ("17", "カ国", "代表の渡航経験"), ("500", "円", "初月お試し価格")]
x_positions = [1.5, 4, 6.5]
for i, (num, unit, label) in enumerate(stats):
    box = add_rounded_box(slide, x_positions[i], 3.2, 2.2, 1.3, RGBColor(44, 82, 130))
    add_textbox(slide, x_positions[i], 3.35, 2.2, 0.7, f"{num}{unit}", 32, True, ACCENT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x_positions[i], 4.0, 2.2, 0.3, label, 11, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0, 5.1, 10, 0.3, "株式会社UNISIA（ユニシア）", 11, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド2: こんなお悩みありませんか？ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 1.0, 0.3, SKY_BLUE)
add_text_frame(tag, "PROBLEM", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "こんなお悩み、ありませんか？", 28, True, PRIMARY, PP_ALIGN.LEFT)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

problems = [
    "😰 航空券、高すぎない？",
    "😩 調べるのが面倒…",
    "😱 家族全員分、いくらになる？",
    "🤔 エージェントは高いって聞く",
    "😕 本当に最安値か分からない",
    "😣 相談できる人がいない"
]

for i, problem in enumerate(problems):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 3.1
    y = 1.5 + row * 1.1
    box = add_rounded_box(slide, x, y, 2.9, 0.9, PINK_LIGHT)
    add_textbox(slide, x + 0.1, y + 0.25, 2.7, 0.5, problem, 13, True, TEXT_DARK, PP_ALIGN.CENTER)

# 解決ボックス
solution = add_rounded_box(slide, 1.5, 4.0, 7, 0.65, SKY_BLUE)
add_text_frame(solution, "これらの悩み、すべて解決します。", 18, True, WHITE, PP_ALIGN.CENTER)

# ===== スライド3: なぜ航空券は高いのか？ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "REASON", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "なぜ、航空券は高いのか？", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# 左側ボックス
left_box = add_rounded_box(slide, 0.4, 1.5, 3.8, 2.0, PINK_LIGHT)
add_textbox(slide, 0.5, 1.7, 3.6, 0.4, "答えはシンプル。", 16, True, DANGER, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.2, 3.6, 0.9, "「探し方を\n知らない」から。", 22, True, DANGER, PP_ALIGN.CENTER)

# 右側の理由
reasons = [
    ("💸", "多くの人は大手比較サイトやエージェントの言い値で買っている"),
    ("🔍", "実はルート・時期・航空会社で数万円〜数十万円変わる"),
    ("⚠️", "エージェントは「最安値」ではなく利益率の高い航空券を提案")
]
y = 1.5
for icon, text in reasons:
    box = add_rounded_box(slide, 4.5, y, 5.1, 0.6, BG_LIGHT)
    add_textbox(slide, 4.6, y + 0.15, 4.9, 0.4, f"{icon} {text}", 11, False, TEXT_MEDIUM)
    y += 0.7

# 下部メッセージ
bottom = add_rounded_box(slide, 0.4, 3.8, 9.2, 0.55, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 3.9, 9, 0.4, "知っている人は、数十万円得している。この差を埋めるのが、フラットサポートです。", 13, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド4: だから私たちが解決します =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 1.0, 0.3, SKY_BLUE)
add_text_frame(tag, "SOLUTION", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "だから、私たちが解決します", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# メインボックス
main_box = add_rounded_box(slide, 0.4, 1.5, 9.2, 1.8, PRIMARY)
add_textbox(slide, 0.5, 1.65, 9, 0.5, "あなたの代わりに、最安値を探します。", 22, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.15, 9, 0.4, "17カ国渡航の経験を持つ代表が、「最安値のルート」を徹底的に探して提案します。", 13, False, WHITE, PP_ALIGN.CENTER)

# ステップ
steps = ["1. LINEで相談", "2. 最安値を探索", "3. 複数プラン提案", "4. 手配完了"]
x = 0.8
for step in steps:
    step_box = add_rounded_box(slide, x, 2.65, 2.0, 0.45, RGBColor(44, 82, 130))
    add_text_frame(step_box, step, 11, True, WHITE, PP_ALIGN.CENTER)
    x += 2.2

# 下部メッセージ
add_textbox(slide, 0.4, 3.6, 9.2, 0.7, "「○月に○○に行きたい」とLINEで送るだけ。\nあとは全部おまかせください。", 15, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド5: フラットサポートとは =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "SERVICE", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "フラットサポートとは", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# 業界初バッジ
badge = add_rounded_box(slide, 4, 1.4, 1.5, 0.35, ACCENT)
add_text_frame(badge, "業界初", 12, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 1.8, 9.2, 0.4, "サブスクリプション型の航空券手配 & 海外挑戦サポートサービス", 16, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 2.2, 9.2, 0.3, "エージェントのように高額なパッケージを売りつけません。", 13, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# フロー
flow_items = ["サブスク登録", "公式LINE追加", "LINEで相談", "サポート開始"]
x = 0.5
for i, item in enumerate(flow_items):
    box = add_rounded_box(slide, x, 2.65, 2.0, 0.5, BG_LIGHT, SKY_BLUE)
    add_text_frame(box, item, 11, True, PRIMARY, PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + 2.05, 2.75, 0.3, 0.3, "→", 18, False, ACCENT, PP_ALIGN.CENTER)
    x += 2.35

# 比較ポイント
add_textbox(slide, 0.4, 3.3, 9.2, 0.3, "【従来のエージェントとの違い】", 13, True, PRIMARY)

comparison = [
    ("料金体系", "高額パッケージ", "✓ シンプル月額制"),
    ("航空券", "利益優先で割高", "✓ 最安値を追求"),
    ("解約", "違約金あり", "✓ いつでも無料")
]

y = 3.65
add_rect(slide, 0.4, y, 2.5, 0.35, PRIMARY)
add_rect(slide, 2.9, y, 3, 0.35, RGBColor(160, 174, 192))
add_rect(slide, 5.9, y, 3.7, 0.35, ACCENT)
add_textbox(slide, 0.5, y + 0.05, 2.3, 0.25, "比較項目", 10, True, WHITE)
add_textbox(slide, 3, y + 0.05, 2.8, 0.25, "従来のエージェント", 10, True, WHITE)
add_textbox(slide, 6, y + 0.05, 3.5, 0.25, "フラットサポート", 10, True, WHITE)

y += 0.4
for item, old, new in comparison:
    add_textbox(slide, 0.5, y, 2.3, 0.28, item, 10, True, TEXT_DARK)
    add_textbox(slide, 3, y, 2.8, 0.28, old, 10, False, TEXT_LIGHT)
    add_textbox(slide, 6, y, 3.5, 0.28, new, 10, False, SUCCESS)
    y += 0.32

# ===== スライド6: 選ばれる3つの理由 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 1.0, 0.3, SKY_BLUE)
add_text_frame(tag, "FEATURES", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "選ばれる3つの理由", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

features = [
    ("💰", "POINT 01", "最大34万円\n安くなる", "大手エージェントより圧倒的に安く航空券を手配。世界一周16万円の実績あり。"),
    ("👨‍👩‍👧‍👦", "POINT 02", "家族全員分\nサポート", "月額8,800円で、ご家族全員の航空券を手配。追加料金なし。最大9名まで。"),
    ("💬", "POINT 03", "LINEで\nいつでも相談", "「この時期はいくら？」「どのルートが安い？」気軽にLINEで何度でも相談OK。")
]

x = 0.4
for icon, point, title, desc in features:
    box = add_rounded_box(slide, x, 1.45, 3.0, 2.3, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.45, 3.0, 0.06, ACCENT)
    add_textbox(slide, x, 1.6, 3.0, 0.5, icon, 32, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.05, 3.0, 0.25, point, 9, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.35, 3.0, 0.6, title, 15, True, PRIMARY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.15, 2.95, 2.7, 0.7, desc, 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    x += 3.2

# 下部メッセージ
bottom = add_rounded_box(slide, 0.4, 3.95, 9.2, 0.6, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 4.05, 9, 0.45, "さらに、初月500円 でお試しいただけます。合わなければ即解約OK。", 14, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド7: 実際の事例 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "RESULTS", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "実際にお得になった事例", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

cases = [
    ("CASE 01", "🌍 世界一周航空券", "50万円", "16万円", "34万円お得！"),
    ("CASE 02", "🇬🇧 イギリス往復", "33万円", "11万円", "22万円お得！"),
    ("CASE 03", "🇦🇺 オーストラリア往復", "18万円", "12万円", "6万円お得！")
]

x = 0.4
for case_num, dest, old_price, new_price, saving in cases:
    box = add_rounded_box(slide, x, 1.45, 3.0, 2.5, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.45, 3.0, 0.06, ACCENT)
    
    badge = add_rounded_box(slide, x + 0.7, 1.6, 1.6, 0.3, PRIMARY)
    add_text_frame(badge, case_num, 10, True, WHITE, PP_ALIGN.CENTER)
    
    add_textbox(slide, x, 2.0, 3.0, 0.35, dest, 14, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.4, 3.0, 0.25, f"大手: {old_price}", 11, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.7, 3.0, 0.35, f"→ {new_price}", 20, True, PRIMARY, PP_ALIGN.CENTER)
    
    saving_box = add_rounded_box(slide, x + 0.3, 3.15, 2.4, 0.45, ACCENT)
    add_text_frame(saving_box, saving, 14, True, WHITE, PP_ALIGN.CENTER)
    
    x += 3.2

add_textbox(slide, 0, 4.1, 10, 0.25, "※時期や条件によって価格は変動します", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# ===== スライド8: 留学費用の比較 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "RESULTS", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "留学費用も大幅削減", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# 左側（大手）
left = add_rounded_box(slide, 0.5, 1.5, 4.2, 1.8, PINK_LIGHT)
add_textbox(slide, 0.6, 1.7, 4.0, 0.25, "大手留学エージェント", 12, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.6, 2.05, 4.0, 0.6, "190万円〜", 36, True, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.6, 2.7, 4.0, 0.25, "6ヶ月留学の場合", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# 右側（フラットサポート）
right = add_rounded_box(slide, 5.0, 1.5, 4.5, 1.8, PRIMARY)
add_textbox(slide, 5.1, 1.7, 4.3, 0.25, "フラットサポート", 12, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 5.1, 2.05, 4.3, 0.6, "65万円〜", 36, True, ACCENT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 5.1, 2.7, 4.3, 0.25, "6ヶ月留学の場合", 10, False, WHITE, PP_ALIGN.CENTER)

# 節約額
saving_box = add_rounded_box(slide, 3.0, 3.5, 4.0, 0.6, ACCENT)
add_text_frame(saving_box, "約125万円お得！", 24, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 4.3, 9.2, 0.5, "同じ留学先・同じ期間でも、取り方次第でこれだけ変わります。\n浮いたお金は、現地での生活費や体験に使えます。", 12, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# ===== スライド9: サポート内容 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "SUPPORT", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "充実のサポート内容", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

add_textbox(slide, 0.4, 1.4, 9.2, 0.3, "航空券手配だけでなく、海外挑戦に必要な10項目のサポートを提供しています。", 12, False, TEXT_MEDIUM)

supports = [
    ("✈️", "航空券の手配", "最安値を徹底追求して手配"),
    ("🚗", "空港送迎手配", "到着後も安心（留学者）"),
    ("📄", "ビザ発行代行", "面倒な手続きを代行"),
    ("🛡️", "海外保険案内", "最適なプランをご提案"),
    ("💬", "24時間LINEサポート", "いつでも何度でも相談OK"),
    ("🆘", "現地トラブル対応", "日本からリモートでサポート"),
    ("🗣️", "現地日本語サポート", "日本からリモートで対応"),
    ("🗺️", "ツアーガイド", "現地でのツアーを開催"),
    ("🎓", "英会話（月4回）", "オンラインマンツーマン"),
    ("💼", "帰国後就職サポート", "求人紹介・履歴書添削・面接対策")
]

for i, (icon, title, desc) in enumerate(supports):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 4.7
    y = 1.75 + row * 0.5
    box = add_rounded_box(slide, x, y, 4.5, 0.45, BG_LIGHT)
    add_textbox(slide, x + 0.1, y + 0.08, 4.3, 0.3, f"{icon} {title} - {desc}", 10, False, TEXT_MEDIUM)

# ===== スライド10: コミュニティ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 1.2, 0.3, SKY_BLUE)
add_text_frame(tag, "COMMUNITY", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "会員限定コミュニティ", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# メインボックス
main = add_rounded_box(slide, 0.4, 1.45, 9.2, 1.2, PRIMARY)
add_textbox(slide, 0.5, 1.55, 9, 0.35, "🎮 Discordコミュニティに無料参加", 18, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.95, 9, 0.6, "フラットサポート会員だけが参加できる限定コミュニティ。\n同じ目標を持つ仲間と繋がり、情報交換や相談ができます。", 13, False, WHITE, PP_ALIGN.CENTER)

# コミュニティの内容
contents = [
    ("👥", "会員同士の交流", "海外経験者・これから行く人が集まり、リアルな情報を共有"),
    ("💡", "先輩からのアドバイス", "留学・ワーホリ経験者が質問に答えてくれる"),
    ("🎉", "イベント開催", "オンライン交流会・オフラインミートアップを定期開催")
]

x = 0.4
for icon, title, desc in contents:
    box = add_rounded_box(slide, x, 2.85, 3.0, 1.25, BG_LIGHT)
    add_textbox(slide, x, 3.0, 3.0, 0.4, icon, 28, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.35, 3.0, 0.3, title, 13, True, PRIMARY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 3.65, 2.8, 0.4, desc, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    x += 3.2

add_textbox(slide, 0.4, 4.25, 9.2, 0.35, "一人で悩まない。仲間がいるから、海外挑戦が楽しくなる。", 14, True, ACCENT, PP_ALIGN.CENTER)

# ===== スライド11: 料金プラン =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "PRICING", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "料金プラン", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# スタンダードプラン
std = add_rounded_box(slide, 0.4, 1.45, 5.5, 3.0, WHITE, ACCENT)
badge = add_rounded_box(slide, 1.8, 1.3, 2.6, 0.35, ACCENT)
add_text_frame(badge, "人気No.1 / 初月500円", 11, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.6, 5.3, 0.35, "スタンダードプラン", 18, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.95, 5.3, 0.25, "航空券手配 + フルサポート", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.2, 5.3, 0.2, "初月", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.35, 5.3, 0.55, "500円", 44, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.85, 5.3, 0.2, "翌月以降 8,800円/月", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

std_features = ["✓ 航空券の最安値検索＆手配", "✓ 家族全員分サポート（追加料金なし）", "✓ 24時間LINEサポート", "✓ Discordコミュニティ参加", "✓ 現地トラブル対応 / 海外保険案内"]
y = 3.1
for feat in std_features:
    add_textbox(slide, 0.7, y, 5.1, 0.22, feat, 10, False, TEXT_MEDIUM)
    y += 0.23

# プレミアムプラン
prem = add_rounded_box(slide, 6.1, 1.45, 3.5, 3.0, WHITE, BG_LIGHT)
add_textbox(slide, 6.2, 1.6, 3.3, 0.35, "プレミアムプラン", 15, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 1.95, 3.3, 0.25, "留学・ワーホリをフルサポート", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 2.35, 3.3, 0.45, "19,800円/月", 24, True, PRIMARY, PP_ALIGN.CENTER)

prem_features = ["✓ スタンダードの内容すべて", "✓ ビザ発行代行", "✓ 月4回オンライン英会話", "✓ 留学フルサポート", "✓ 帰国後就職サポート"]
y = 2.9
for feat in prem_features:
    add_textbox(slide, 6.3, y, 3.1, 0.22, feat, 9, False, TEXT_MEDIUM)
    y += 0.23

# 下部メッセージ
bottom = add_rounded_box(slide, 0.4, 4.55, 9.2, 0.45, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 4.6, 9, 0.35, "いつでも解約OK・解約手数料なし", 13, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド12: 他社比較 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.9, 0.3, SKY_BLUE)
add_text_frame(tag, "COMPARE", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "留学エージェントとの違い", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# ヘッダー
add_rect(slide, 0.4, 1.45, 2.6, 0.4, PRIMARY)
add_rect(slide, 3.0, 1.45, 3.2, 0.4, RGBColor(160, 174, 192))
add_rect(slide, 6.2, 1.45, 3.4, 0.4, ACCENT)
add_textbox(slide, 0.5, 1.5, 2.4, 0.3, "比較項目", 11, True, WHITE)
add_textbox(slide, 3.1, 1.5, 3.0, 0.3, "留学エージェント", 11, True, WHITE)
add_textbox(slide, 6.3, 1.5, 3.2, 0.3, "フラットサポート", 11, True, WHITE)

comparison = [
    ("対象", "留学する人のみ", "✓ 旅行でも留学でもOK"),
    ("料金体系", "パッケージ（高額）", "✓ シンプルな月額制"),
    ("航空券", "割高なことが多い", "✓ 最安値を徹底追求"),
    ("家族利用", "追加料金あり", "✓ 追加料金なし（最大9名）"),
    ("相談方法", "予約制の面談", "✓ LINEでいつでも"),
    ("解約", "違約金あり", "✓ いつでも無料"),
    ("お試し", "なし", "✓ 初月500円"),
    ("コミュニティ", "なし or 有料", "✓ 無料で参加可能")
]

y = 1.9
for item, old, new in comparison:
    add_textbox(slide, 0.5, y, 2.4, 0.28, item, 10, True, TEXT_DARK)
    add_textbox(slide, 3.1, y, 3.0, 0.28, old, 10, False, TEXT_LIGHT)
    add_textbox(slide, 6.3, y, 3.2, 0.28, new, 10, False, SUCCESS)
    y += 0.32

# ===== スライド13: お客様の声1 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.7, 0.3, SKY_BLUE)
add_text_frame(tag, "VOICE", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "ご利用者様の声", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

voices = [
    (True, "Y", "Y.Tさん（20代・男性）", "世界一周で利用", "「世界一周の航空券が合計16万円で取れました！大手だと50万円近くかかると言われていたので、正直信じられなかったです。」"),
    (False, "M", "M.Sさん（20代・女性）", "ワーホリ準備で利用", "「航空券だけで6万円浮いて、その分を現地での生活費に回せました。余計なパッケージを押し付けられないのも良かったです！」"),
    (False, "K", "K.Nさん（30代・男性）", "家族旅行で利用", "「家族4人でハワイに行きました。自分で調べた価格より8万円も安く取れました。子どもがいると調べる時間もないので、LINEで全部やってくれるのが助かりました！」")
]

y = 1.45
for featured, initial, name, tag_text, text in voices:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.95, PRIMARY if featured else BG_LIGHT)
    text_color = WHITE if featured else TEXT_MEDIUM
    name_color = WHITE if featured else TEXT_DARK
    tag_color = ACCENT_LIGHT if featured else TEXT_LIGHT
    
    add_textbox(slide, 0.6, y + 0.1, 3, 0.25, f"{initial} | {name}", 11, True, name_color)
    add_textbox(slide, 0.6, y + 0.35, 1.5, 0.2, tag_text, 9, False, tag_color)
    add_textbox(slide, 0.6, y + 0.55, 8.8, 0.35, text, 10, False, text_color)
    y += 1.05

# ===== スライド14: お客様の声2 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.7, 0.3, SKY_BLUE)
add_text_frame(tag, "VOICE", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "ご利用者様の声", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

voices2 = [
    (False, "S", "S.Hさん（30代・女性）", "イギリス留学で利用", "「イギリス留学の航空券、大手だと33万円と言われたのが11万円で取れました。22万円も浮いたおかげで、現地での語学学校の期間を延ばせました。」"),
    (True, "R", "R.Kさん（20代・女性）", "LINEサポートが決め手", "「何よりLINEでのサポートが手厚いのが良かったです！『この時期に行きたいんだけど…』って相談したら、すぐに複数のルートを提案してくれました。」")
]

y = 1.45
for featured, initial, name, tag_text, text in voices2:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.95, PRIMARY if featured else BG_LIGHT)
    text_color = WHITE if featured else TEXT_MEDIUM
    name_color = WHITE if featured else TEXT_DARK
    tag_color = ACCENT_LIGHT if featured else TEXT_LIGHT
    
    add_textbox(slide, 0.6, y + 0.1, 3, 0.25, f"{initial} | {name}", 11, True, name_color)
    add_textbox(slide, 0.6, y + 0.35, 1.5, 0.2, tag_text, 9, False, tag_color)
    add_textbox(slide, 0.6, y + 0.55, 8.8, 0.35, text, 10, False, text_color)
    y += 1.05

# 下部メッセージ
msg = add_rounded_box(slide, 0.4, 3.7, 9.2, 0.55, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 3.8, 9, 0.4, "多くの方に選ばれています。次はあなたが、お得に海外へ行く番です。", 14, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド15: 代表紹介 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.7, 0.3, SKY_BLUE)
add_text_frame(tag, "ABOUT", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "代表紹介", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# 左側
add_textbox(slide, 0.4, 1.5, 3.5, 0.35, "🌍 17カ国渡航", 14, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 1.9, 3.5, 0.45, "井上 智羅（トライ）", 20, True, TEXT_DARK, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 2.35, 3.5, 0.25, "株式会社UNISIA 代表取締役", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# 右側ストーリー
story = """2020年、専門学校卒業後にマルタ・NZ留学へ。
その後、世界一周の旅に出てノマドフリーランスとして活動。

一方で、留学エージェントに200万円損した経験も。
高額なパッケージを売りつけられ、後から
「もっと安くできた」と知りました。"""

add_textbox(slide, 4.0, 1.5, 5.5, 1.5, story, 11, False, TEXT_MEDIUM)

# メッセージボックス
msg_box = add_rounded_box(slide, 0.4, 3.3, 9.2, 1.0, BG_LIGHT)
add_textbox(slide, 0.5, 3.45, 9, 0.8, "「あの時の自分と同じ思いをしてほしくない。」\nだから、フラットサポートを作りました。あなたの海外挑戦を、全力でサポートします。", 13, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド16: よくある質問1 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.6, 0.3, SKY_BLUE)
add_text_frame(tag, "FAQ", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "よくある質問", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

faqs = [
    ("Q. 本当に安くなりますか？", "A. はい。大手エージェントと比較して、最大34万円お得になった実績があります。常に最安値を徹底的にお探しします。"),
    ("Q. 家族は何人まで使えますか？", "A. 1契約につき最大9名様までご利用いただけます。ご家族全員分を追加料金なしで手配します。"),
    ("Q. 解約はいつでもできますか？", "A. はい。いつでも解約可能です。解約手数料もかかりません。初月500円で合わなければそのまま解約できます。")
]

y = 1.45
for q, a in faqs:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.85, BG_LIGHT)
    add_textbox(slide, 0.6, y + 0.1, 8.8, 0.25, q, 13, True, TEXT_DARK)
    add_textbox(slide, 0.6, y + 0.4, 8.8, 0.4, a, 10, False, TEXT_MEDIUM)
    y += 0.95

# ===== スライド17: よくある質問2 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])

tag = add_rounded_box(slide, 0.4, 0.3, 0.6, 0.3, SKY_BLUE)
add_text_frame(tag, "FAQ", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 0.7, 9, 0.6, "よくある質問", 28, True, PRIMARY)
add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

faqs2 = [
    ("Q. 旅行だけでも使えますか？留学じゃなくても大丈夫？", "A. もちろんです。留学だけでなく、海外旅行、ワーホリ、ビジネス出張など、あらゆる海外渡航に対応しています。"),
    ("Q. コミュニティはどんな内容ですか？", "A. Discordで会員限定コミュニティを運営しています。会員同士の交流、先輩からのアドバイス、オンライン/オフラインイベントなどがあります。"),
    ("Q. 英会話はどんな形式ですか？", "A. プレミアムプラン限定で、月4回のオンラインマンツーマンレッスンを提供しています。")
]

y = 1.45
for q, a in faqs2:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.85, BG_LIGHT)
    add_textbox(slide, 0.6, y + 0.1, 8.8, 0.25, q, 13, True, TEXT_DARK)
    add_textbox(slide, 0.6, y + 0.4, 8.8, 0.4, a, 10, False, TEXT_MEDIUM)
    y += 0.95

# ===== スライド18: まとめ・CTA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.6, 10, 0.7, "まずは初月500円で\nお試しください", 32, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.5, 10, 0.35, "航空券1回の手配で、数万円〜数十万円の差が出ます。", 15, False, WHITE, PP_ALIGN.CENTER)

# 価格ボックス
price_box = add_rounded_box(slide, 3.2, 2.0, 3.6, 1.2, RGBColor(44, 82, 130))
add_textbox(slide, 3.3, 2.1, 3.4, 0.25, "初月", 12, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 3.3, 2.35, 3.4, 0.7, "500円", 56, True, ACCENT_LIGHT, PP_ALIGN.CENTER)

# 特徴
add_textbox(slide, 0, 3.4, 10, 0.35, "✓ 家族全員分OK　　✓ LINEでいつでも相談　　✓ いつでも解約可能　　✓ コミュニティ参加", 13, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0, 3.9, 10, 0.25, "※ 翌月以降 8,800円/月 ・ 解約手数料なし", 11, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド19: お問い合わせ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.8, 10, 0.6, "お問い合わせ", 36, True, WHITE, PP_ALIGN.CENTER)

# 連絡先ボックス
contact_box = add_rounded_box(slide, 2.2, 1.7, 5.6, 2.2, RGBColor(44, 82, 130))

contacts = [
    ("会社名", "株式会社UNISIA（ユニシア）"),
    ("代表取締役", "井上 智羅（いのうえ とらい）"),
    ("所在地", "福岡県福岡市博多区博多駅前1丁目23番2号"),
    ("サービスURL", "ltdunisia.memberpay.jp"),
    ("公式LINE", "@unisia")
]

y = 1.85
for label, value in contacts:
    add_textbox(slide, 2.4, y, 1.8, 0.28, label, 11, False, WHITE)
    add_textbox(slide, 4.3, y, 3.3, 0.28, value, 11, False, WHITE)
    y += 0.38

add_textbox(slide, 0, 4.2, 10, 0.4, '"すべての人に、気軽な海外挑戦を"', 18, False, WHITE, PP_ALIGN.CENTER)

# 保存
output_path = "/Users/user/Library/Mobile Documents/com~apple~CloudDocs/Odsidian/トライ/02_ビジネス/ユニシア/projects/sales-material/フラットサポート営業資料_v2.pptx"
prs.save(output_path)
print(f"PowerPointファイルを作成しました: {output_path}")
