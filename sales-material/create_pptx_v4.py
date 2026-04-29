#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
フラットサポート営業資料 PowerPoint作成スクリプト v4
- 文字の被り修正
- サポート内容を詳しく説明
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラー定義
PRIMARY = RGBColor(26, 54, 93)
ACCENT = RGBColor(237, 137, 54)
ACCENT_LIGHT = RGBColor(246, 173, 85)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(26, 32, 44)
TEXT_MEDIUM = RGBColor(74, 85, 104)
TEXT_LIGHT = RGBColor(113, 128, 150)
SKY_BLUE = RGBColor(66, 153, 225)
SUCCESS = RGBColor(72, 187, 120)
DANGER = RGBColor(229, 62, 62)
BG_LIGHT = RGBColor(247, 250, 252)
BG_CREAM = RGBColor(255, 250, 240)
PINK_LIGHT = RGBColor(254, 215, 215)

def add_text_frame(shape, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tf

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_rounded_box(slide, left, top, width, height, fill_color, line_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if line_color:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box

def add_rect(slide, left, top, width, height, fill_color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    rect.line.fill.background()
    return rect

def add_section_header(slide, tag_text, title):
    tag = add_rounded_box(slide, 0.4, 0.25, len(tag_text) * 0.11 + 0.4, 0.28, SKY_BLUE)
    add_text_frame(tag, tag_text, 9, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, 0.4, 0.6, 9, 0.5, title, 26, True, PRIMARY)
    add_rect(slide, 0.4, 1.05, 9.2, 0.04, ACCENT)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ===== スライド1: 表紙 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.25, 10, 0.35, "初月500円キャンペーン中", 13, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.7, 10, 0.25, "LTD. UNISIA", 11, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.4, 10, 0.7, "フラットサポート", 44, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.2, 10, 0.4, "月額制の航空券手配 & 海外挑戦サポート", 17, False, WHITE, PP_ALIGN.CENTER)

stats = [("34", "万円", "最大節約額"), ("17", "カ国", "代表の渡航経験"), ("500", "円", "初月お試し価格")]
x_positions = [1.5, 4, 6.5]
for i, (num, unit, label) in enumerate(stats):
    box = add_rounded_box(slide, x_positions[i], 2.9, 2.2, 1.2, RGBColor(44, 82, 130))
    add_textbox(slide, x_positions[i], 3.0, 2.2, 0.6, f"{num}{unit}", 30, True, ACCENT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x_positions[i], 3.6, 2.2, 0.3, label, 10, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0, 4.7, 10, 0.25, "株式会社UNISIA（ユニシア）", 10, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド2: こんなお悩みありませんか？ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "PROBLEM", "こんなお悩み、ありませんか？")

problems = [
    ("😰", "航空券、高すぎない？"),
    ("😩", "調べるのが面倒…"),
    ("😱", "家族全員分、いくらになる？"),
    ("🤔", "エージェントは高いって聞く"),
    ("😕", "本当に最安値か分からない"),
    ("😣", "相談できる人がいない")
]

for i, (icon, text) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 3.1
    y = 1.25 + row * 1.0
    box = add_rounded_box(slide, x, y, 2.9, 0.85, PINK_LIGHT)
    add_textbox(slide, x, y + 0.15, 2.9, 0.25, icon, 20, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 0.45, 2.9, 0.3, text, 12, True, TEXT_DARK, PP_ALIGN.CENTER)

solution = add_rounded_box(slide, 1.5, 3.55, 7, 0.6, SKY_BLUE)
add_text_frame(solution, "これらの悩み、すべて解決します。", 16, True, WHITE, PP_ALIGN.CENTER)

# ===== スライド3: フラットサポートとは =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SERVICE", "フラットサポートとは")

problem_box = add_rounded_box(slide, 0.4, 1.2, 4.4, 1.2, PINK_LIGHT)
add_textbox(slide, 0.5, 1.3, 4.2, 0.25, "💸 なぜ航空券は高い？", 13, True, DANGER, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.6, 4.2, 0.65, "多くの人はエージェントの言い値で購入。\n実は取り方次第で数万円〜数十万円変わる。", 10, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

solution_box = add_rounded_box(slide, 5.0, 1.2, 4.6, 1.2, PRIMARY)
add_textbox(slide, 5.1, 1.3, 4.4, 0.25, "✈️ だから私たちが探します", 13, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 5.1, 1.6, 4.4, 0.65, "17カ国渡航の代表が最安値を徹底追求。\nLINEで相談するだけで全部おまかせ。", 10, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 2.55, 9.2, 0.3, "業界初のサブスクリプション型 航空券手配 & 海外挑戦サポートサービス", 13, True, ACCENT, PP_ALIGN.CENTER)

flow_items = ["サブスク登録", "公式LINE追加", "LINEで相談", "サポート開始"]
x = 0.5
for i, item in enumerate(flow_items):
    box = add_rounded_box(slide, x, 2.95, 2.0, 0.4, BG_LIGHT, SKY_BLUE)
    add_text_frame(box, item, 10, True, PRIMARY, PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + 2.05, 3.0, 0.3, 0.3, "→", 14, False, ACCENT, PP_ALIGN.CENTER)
    x += 2.35

comparison = [
    ("✓ 料金: 高額パッケージ → シンプル月額制"),
    ("✓ 航空券: 割高 → 最安値追求（スカイスキャナー・trip.comより安く）"),
    ("✓ 解約: 違約金あり → いつでも無料")
]
y = 3.5
for item in comparison:
    add_textbox(slide, 0.5, y, 9, 0.22, item, 10, False, TEXT_MEDIUM)
    y += 0.25

# ===== スライド4: 選ばれる3つの理由 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "FEATURES", "選ばれる3つの理由")

features = [
    ("💰", "POINT 01", "最大34万円安くなる", "大手エージェントより圧倒的に安く。\n世界一周16万円の実績。"),
    ("👨‍👩‍👧‍👦", "POINT 02", "家族全員分サポート", "月額8,800円で最大9名まで。\n追加料金なし。"),
    ("💬", "POINT 03", "LINEでいつでも相談", "「この時期いくら？」\n気軽に何度でもOK。")
]

x = 0.4
for icon, point, title, desc in features:
    box = add_rounded_box(slide, x, 1.2, 3.0, 1.9, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.2, 3.0, 0.05, ACCENT)
    add_textbox(slide, x, 1.35, 3.0, 0.35, icon, 26, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 1.7, 3.0, 0.2, point, 9, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 1.95, 3.0, 0.3, title, 13, True, PRIMARY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.3, 2.8, 0.65, desc, 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    x += 3.2

bottom = add_rounded_box(slide, 0.4, 3.25, 9.2, 0.5, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 3.35, 9, 0.35, "さらに、初月500円 でお試し。合わなければ即解約OK。", 13, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド5: 実際の事例 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "RESULTS", "実際にお得になった事例")

cases = [
    ("CASE 01", "🌍 世界一周航空券", "50万円", "16万円", "34万円お得！"),
    ("CASE 02", "🇬🇧 イギリス往復", "33万円", "11万円", "22万円お得！"),
    ("CASE 03", "🇦🇺 オーストラリア往復", "18万円", "12万円", "6万円お得！")
]

x = 0.4
for case_num, dest, old_price, new_price, saving in cases:
    box = add_rounded_box(slide, x, 1.2, 3.0, 2.0, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.2, 3.0, 0.05, ACCENT)
    badge = add_rounded_box(slide, x + 0.7, 1.3, 1.6, 0.25, PRIMARY)
    add_text_frame(badge, case_num, 9, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x, 1.6, 3.0, 0.25, dest, 12, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 1.9, 3.0, 0.2, f"大手: {old_price}", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.15, 3.0, 0.3, f"→ {new_price}", 17, True, PRIMARY, PP_ALIGN.CENTER)
    saving_box = add_rounded_box(slide, x + 0.35, 2.5, 2.3, 0.35, ACCENT)
    add_text_frame(saving_box, saving, 12, True, WHITE, PP_ALIGN.CENTER)
    x += 3.2

add_textbox(slide, 0.4, 3.4, 9.2, 0.3, "📚 6ヶ月留学: 大手190万円 → フラットサポート65万円 = 約125万円お得！", 12, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.75, 10, 0.2, "※時期や条件によって価格は変動します", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# ===== スライド6: サポート内容①（航空券手配） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 01", "航空券の手配")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "✈️ スカイスキャナー・trip.comより安い「本当の最安値」を提示", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("🔍", "徹底比較", "複数の航空会社・ルート・乗り継ぎ便を全て比較。格安サイトでは見つからない隠れた最安値ルートを発見します。"),
    ("📅", "最適タイミング", "航空券は購入時期で価格が大きく変動。過去データを元に、最もお得な購入タイミングをご提案します。"),
    ("💰", "購入サポート", "最安値の提示だけでなく、実際に購入できる状態まで丁寧にサポート。面倒な手続きもおまかせ。"),
    ("🌐", "国内・国外対応", "海外航空券だけでなく、国内航空券の手配も対応。出張・帰省など幅広いニーズに対応します。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.5, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.6, y + 0.08, 6.8, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド7: サポート内容②（海外保険） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 02", "海外保険案内")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "🛡️ やり方次第で保険料0円も可能！最適なプランをご提案", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("💳", "クレカ活用", "クレジットカード付帯保険を最大限活用。複数カードの組み合わせで、追加保険なしでも十分な補償が得られるケースも。"),
    ("📊", "比較提案", "渡航先・期間・目的に応じた保険を徹底比較。長期滞在向けの割安プランや、必要な補償だけに絞ったプランもご紹介。"),
    ("🆘", "請求サポート", "万が一の際の保険金請求もサポート。必要書類の案内から手続きの流れまで、日本語で丁寧に対応します。"),
    ("💡", "0円プラン相談", "「保険料をできるだけ抑えたい」というご要望にも対応。条件次第で保険料0円で渡航できる方法もご案内します。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.5, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.6, y + 0.08, 6.8, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド8: サポート内容③（現地日本語対応） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 03", "現地日本語対応サポート")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "🗣️ 留学生向け：現地で問題が起きても日本人がしっかり対応", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("🏫", "学校トラブル対応", "入学手続きのトラブル、クラス変更の相談、学校との交渉など、現地の日本人スタッフが直接対応します。"),
    ("🏥", "緊急時サポート", "病気・ケガ・事故など緊急時も安心。病院の付き添いや、現地での各種手続きを日本人がサポート。"),
    ("🏠", "生活サポート", "銀行口座開設、携帯電話契約、住居トラブルなど、現地生活で困ったことは何でも相談できます。"),
    ("🤝", "安心の日本人対応", "言葉の壁で困ることなく、日本語でしっかりコミュニケーション。留学生活を安心してスタートできます。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.8, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.08, 6.5, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド9: サポート内容④（LINEサポート＆情勢配信） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 04", "24時間LINEサポート & 海外情勢配信")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "💬 いつでも日本語で相談OK + リアルタイム海外情報を配信", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("📱", "24時間LINE対応", "時差があっても安心。現地で困ったこと、航空券の相談、些細な質問まで、LINEでいつでも日本語で対応します。"),
    ("🌍", "海外情勢配信", "各国の治安情報、政治情勢、天気などをリアルタイムで配信。「今ウクライナは安全？」「アメリカの天気は？」など最新情報をお届け。"),
    ("⚠️", "渡航注意情報", "外務省の危険情報、テロ・デモ情報、自然災害など、渡航に影響する情報を素早くお知らせ。安全な旅行計画をサポート。"),
    ("💡", "旅行アドバイス", "おすすめスポット、現地のお得情報、トラブル回避のコツなど、経験豊富なスタッフが実践的なアドバイスを提供します。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.8, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.08, 6.5, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド10: サポート内容⑤（留学サポート） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 05", "留学サポート")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "🎓 あなたのニーズに合った学校を提案", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("🗺️", "旅行重視タイプ", "「留学中に旅行もしたい！」という方には、休暇が取りやすい・立地の良い学校をご提案。観光と学習の両立が可能。"),
    ("🎭", "文化体験タイプ", "「現地の文化を楽しみたい！」という方には、アクティビティが充実・現地交流が多い学校をご提案。"),
    ("📚", "語学集中タイプ", "「英語力を本気で伸ばしたい！」という方には、授業時間が長い・厳しめのカリキュラムの学校をご提案。"),
    ("🏠", "滞在先サポート", "ホームステイ・学生寮・シェアハウスなど、ご希望に合った滞在先の手配もサポートします。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.8, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.08, 6.5, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド11: サポート内容⑥（ビザ代行） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 06", "ビザ発行代行")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "📄 面倒なビザ申請をサポート → 今後は丸投げ対応も予定", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("📋", "申請サポート", "必要書類の案内から申請方法まで丁寧にサポート。初めてのビザ申請でも安心して進められます。"),
    ("⚠️", "エラー回避", "特にインドなど、ビザ申請時にエラーが頻発する国も対応。トラブルを未然に防ぎ、スムーズな発行をサポート。"),
    ("🔄", "丸投げ対応（予定）", "今後は完全代行サービスを予定。書類準備から申請まで全てお任せいただけるようになります。"),
    ("🌏", "各国対応", "アメリカ・イギリス・オーストラリア・カナダ・インドなど、主要国のビザ申請に対応しています。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.8, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.08, 6.5, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド12: サポート内容⑦（英会話教室） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 07", "英会話教室（月4回）")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "🎓 TOEIC800点以上の講師による「話せる」ための英会話", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("🗣️", "話せるを目指す", "「聞き取れる」「分かる」ではなく「話せる」状態を目指します。実践的な会話トレーニングで即戦力に。"),
    ("🏫", "入学テスト対策", "語学学校の入学テストで上位クラスからスタートできるよう、事前に英語力を底上げします。"),
    ("✈️", "渡航準備", "現地で困らないレベルの日常会話を習得。空港・ホテル・レストランなど実践的なシーンを想定した練習。"),
    ("👨‍🏫", "マンツーマン", "オンラインでのマンツーマンレッスン。あなたのペースに合わせた個別指導で効率的に上達。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.8, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.08, 6.5, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

add_textbox(slide, 0.4, 4.4, 9.2, 0.2, "※プレミアムプラン限定のサービスです", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# ===== スライド13: サポート内容⑧（帰国後就職） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 08", "帰国後就職サポート")

main_box = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.7, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.5, "💼 海外経験を活かしたキャリア支援", 15, True, WHITE, PP_ALIGN.CENTER)

points = [
    ("🔍", "求人紹介", "海外経験を活かせる企業とのマッチング。語学力や異文化経験を評価してくれる求人をご紹介します。"),
    ("📝", "書類添削", "履歴書・職務経歴書の添削。海外経験をどうアピールするか、採用担当者に響く書き方をアドバイス。"),
    ("🎤", "面接対策", "海外経験のアピール方法、よくある質問への回答準備など、面接で自信を持てるようサポート。"),
    ("📈", "キャリア相談", "「海外経験を活かして何ができる？」将来のキャリアプラン設計から、具体的な就職活動までサポートします。")
]

y = 1.95
for icon, title, desc in points:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.55, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.08, 0.4, 0.35, icon, 18, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.08, 1.5, 0.2, title, 11, True, PRIMARY)
    add_textbox(slide, 2.6, y + 0.08, 6.8, 0.4, desc, 9, False, TEXT_MEDIUM)
    y += 0.6

# ===== スライド14: コミュニティ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "COMMUNITY", "会員限定コミュニティ")

main = add_rounded_box(slide, 0.4, 1.15, 9.2, 0.6, PRIMARY)
add_textbox(slide, 0.5, 1.25, 9, 0.4, "🎮 2つのDiscordコミュニティに無料参加", 15, True, WHITE, PP_ALIGN.CENTER)

# 2つのコミュニティ
comm1 = add_rounded_box(slide, 0.4, 1.85, 4.5, 1.1, BG_LIGHT, SKY_BLUE)
add_textbox(slide, 0.5, 1.95, 4.3, 0.25, "🌍 海外旅行好きコミュニティ", 12, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.25, 4.3, 0.6, "海外旅行が好きな人が集まり、\nおすすめスポット・安い航空券情報・\n現地のリアルな体験談を共有", 9, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

comm2 = add_rounded_box(slide, 5.1, 1.85, 4.5, 1.1, BG_LIGHT, ACCENT)
add_textbox(slide, 5.2, 1.95, 4.3, 0.25, "🗾 国内旅行好きコミュニティ", 12, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 5.2, 2.25, 4.3, 0.6, "国内旅行が好きな人が集まり、\n穴場スポット・安く泊まれる宿・\n地域のおすすめ情報を共有", 9, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# メリット
add_textbox(slide, 0.4, 3.1, 9.2, 0.25, "【コミュニティのメリット】", 11, True, PRIMARY)
merits = [
    "✓ 自分の知らない地域・スポットを知ることができる",
    "✓ おすすめの宿・安く泊まれる場所を教えてもらえる",
    "✓ 先輩会員からリアルな体験談・アドバイスがもらえる",
    "✓ 旅行仲間・一緒に海外に行く仲間が見つかる"
]
y = 3.4
for merit in merits:
    add_textbox(slide, 0.5, y, 9, 0.2, merit, 10, False, TEXT_MEDIUM)
    y += 0.23

# ===== スライド15: その他サポート =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT", "その他のサポート")

supports = [
    ("🚗", "空港送迎手配", "現地空港から滞在先までの送迎を手配。深夜・早朝便でも対応。日本語対応ドライバーも手配可能。"),
    ("🗺️", "ツアーガイド", "会員限定の現地ツアーを定期開催。観光名所だけでなくローカルスポットも案内。他会員との交流の機会にも。"),
    ("📷", "海外現地撮影", "プロカメラマンによる撮影サービス。留学・ワーホリの記念撮影、SNS映えするスポットでの撮影など。"),
    ("🆘", "現地トラブル対応", "パスポート紛失・盗難、航空便欠航時の振替手配など、渡航中のトラブルに日本からリモートで対応。"),
    ("📊", "為替・費用相談", "現地での両替のタイミング、クレジットカードの使い方、現地での生活費の目安など、お金に関する相談も対応。"),
    ("🎁", "会員限定特典", "提携先のホテル・アクティビティの割引、会員限定のお得情報など、様々な特典をご用意。")
]

y = 1.15
for icon, title, desc in supports:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.48, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.06, 0.4, 0.3, icon, 16, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, y + 0.06, 1.8, 0.18, title, 11, True, PRIMARY)
    add_textbox(slide, 2.9, y + 0.06, 6.5, 0.35, desc, 9, False, TEXT_MEDIUM)
    y += 0.52

# ===== スライド16: 料金プラン =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "PRICING", "料金プラン")

# スタンダードプラン
std = add_rounded_box(slide, 0.4, 1.15, 4.4, 2.65, WHITE, ACCENT)
badge = add_rounded_box(slide, 1.2, 1.0, 2.4, 0.28, ACCENT)
add_text_frame(badge, "人気No.1 / 初月500円", 9, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.22, 4.2, 0.26, "スタンダードプラン", 15, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.5, 4.2, 0.16, "航空券手配 + フルサポート", 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.68, 4.2, 0.14, "初月", 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.8, 4.2, 0.4, "500円", 34, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.18, 4.2, 0.14, "翌月以降 8,800円/月", 8, False, TEXT_LIGHT, PP_ALIGN.CENTER)

std_features = ["✓ 航空券の最安値検索＆手配", "✓ 家族全員分サポート（追加料金なし）", "✓ 24時間LINEサポート＆海外情勢配信", "✓ Discordコミュニティ（2種類）参加", "✓ 現地トラブル対応 / 海外保険案内", "✓ 帰国後就職サポート"]
y = 2.38
for feat in std_features:
    add_textbox(slide, 0.55, y, 4.1, 0.16, feat, 8, False, TEXT_MEDIUM)
    y += 0.17

# プレミアムプラン（ディープブルー・特別感演出）
PREMIUM_BLUE = RGBColor(20, 40, 80)
PREMIUM_ACCENT = RGBColor(100, 180, 255)

prem = add_rounded_box(slide, 4.95, 1.15, 4.7, 2.65, PREMIUM_BLUE)
prem_badge = add_rounded_box(slide, 6.5, 0.95, 1.9, 0.32, DANGER)
add_text_frame(prem_badge, "30名限定", 10, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 5.05, 1.22, 4.5, 0.26, "プレミアムプラン", 15, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 5.05, 1.5, 4.5, 0.16, "留学・ワーホリをフルサポート", 8, False, PREMIUM_ACCENT, PP_ALIGN.CENTER)

# 価格表示（お得感演出）
add_textbox(slide, 5.05, 1.68, 4.5, 0.12, "通常50,000円相当が", 8, False, RGBColor(180, 180, 180), PP_ALIGN.CENTER)
add_textbox(slide, 5.05, 1.78, 4.5, 0.4, "19,800円/月", 26, True, PREMIUM_ACCENT, PP_ALIGN.CENTER)

# 含まれるサービスの価値内訳
add_textbox(slide, 5.1, 2.18, 4.4, 0.14, "【含まれるサービスの価値】", 8, True, WHITE)
prem_values = [
    ("専属サポート担当", "10,000円相当"),
    ("優先対応", "5,000円相当"),
    ("月4回オンライン英会話", "15,000円相当"),
    ("ビザ発行代行", "10,000円相当"),
    ("空港送迎手配", "5,000円相当"),
    ("留学フルサポート", "5,000円相当")
]
y = 2.34
for service, value in prem_values:
    add_textbox(slide, 5.15, y, 2.6, 0.13, f"✓ {service}", 7, False, WHITE)
    add_textbox(slide, 7.8, y, 1.7, 0.13, value, 7, False, PREMIUM_ACCENT, PP_ALIGN.RIGHT)
    y += 0.14

# 全部込みアピール
all_in_box = add_rounded_box(slide, 5.15, 3.2, 4.3, 0.35, PREMIUM_ACCENT)
add_textbox(slide, 5.25, 3.25, 4.1, 0.25, "これ全部込みで月額19,800円", 10, True, PREMIUM_BLUE, PP_ALIGN.CENTER)

bottom = add_rounded_box(slide, 0.4, 3.9, 9.2, 0.35, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 3.95, 9, 0.25, "いつでも解約OK・解約手数料なし", 11, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド17: 他社比較 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "COMPARE", "留学エージェントとの違い")

add_rect(slide, 0.4, 1.15, 2.6, 0.35, PRIMARY)
add_rect(slide, 3.0, 1.15, 3.2, 0.35, RGBColor(160, 174, 192))
add_rect(slide, 6.2, 1.15, 3.4, 0.35, ACCENT)
add_textbox(slide, 0.5, 1.2, 2.4, 0.25, "比較項目", 10, True, WHITE)
add_textbox(slide, 3.1, 1.2, 3.0, 0.25, "留学エージェント", 10, True, WHITE)
add_textbox(slide, 6.3, 1.2, 3.2, 0.25, "フラットサポート", 10, True, WHITE)

comparison = [
    ("対象", "留学する人のみ", "✓ 旅行でも留学でもOK"),
    ("料金体系", "パッケージ（高額）", "✓ シンプルな月額制"),
    ("航空券", "割高なことが多い", "✓ 最安値を徹底追求"),
    ("家族利用", "追加料金あり", "✓ 追加料金なし（最大9名）"),
    ("相談方法", "予約制の面談", "✓ LINEでいつでも"),
    ("解約", "違約金あり", "✓ いつでも無料"),
    ("お試し", "なし", "✓ 初月500円"),
    ("コミュニティ", "なし or 有料", "✓ 無料で2種類参加可能")
]

y = 1.55
for item, old, new in comparison:
    add_textbox(slide, 0.5, y, 2.4, 0.22, item, 9, True, TEXT_DARK)
    add_textbox(slide, 3.1, y, 3.0, 0.22, old, 9, False, TEXT_LIGHT)
    add_textbox(slide, 6.3, y, 3.2, 0.22, new, 9, False, SUCCESS)
    y += 0.27

# ===== スライド18: お客様の声 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "VOICE", "ご利用者様の声")

voices = [
    (True, "Y.T", "20代男性｜世界一周", "「世界一周の航空券が合計16万円で取れました！大手だと50万円近くかかると言われていたので、正直信じられなかったです。」"),
    (False, "M.S", "20代女性｜ワーホリ", "「航空券だけで6万円浮いて、その分を現地での生活費に回せました。パッケージを押し付けられないのが良かったです！」"),
    (False, "K.N", "30代男性｜家族旅行", "「家族4人でハワイに行きました。自分で調べた価格より8万円も安く。LINEで全部やってくれるのが助かりました！」"),
    (True, "S.H", "30代女性｜イギリス留学", "「大手だと33万円が11万円に。22万円も浮いて語学学校の期間を延ばせました。」")
]

y = 1.15
for featured, name, tag, text in voices:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.62, PRIMARY if featured else BG_LIGHT)
    text_color = WHITE if featured else TEXT_MEDIUM
    add_textbox(slide, 0.55, y + 0.06, 2.2, 0.22, f"{name}｜{tag}", 9, True, WHITE if featured else TEXT_DARK)
    add_textbox(slide, 0.55, y + 0.3, 8.9, 0.28, text, 9, False, text_color)
    y += 0.68

# ===== スライド19: 代表紹介 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "ABOUT", "代表紹介")

add_textbox(slide, 0.4, 1.2, 3.2, 0.25, "🌍 17カ国渡航", 13, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 1.5, 3.2, 0.35, "井上 智羅（トライ）", 17, True, TEXT_DARK, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 1.85, 3.2, 0.2, "株式会社UNISIA 代表取締役", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)

career = "2020年 専門学校卒業後、マルタ・NZ留学\n2021年 世界一周の旅 & ノマドフリーランス\n2021年 株式会社ファーストビュー提携\n　　　  （ウェブスキ設立）\n2025年 株式会社Unisia設立"
add_textbox(slide, 3.8, 1.2, 5.8, 0.95, career, 9, False, TEXT_MEDIUM)

story_box = add_rounded_box(slide, 0.4, 2.3, 9.2, 1.4, BG_LIGHT)
add_textbox(slide, 0.5, 2.4, 9, 1.2, 
"留学エージェントに200万円損した経験があります。高額なパッケージを売りつけられ、後から「もっと安くできた」と知りました。\n\n「あの時の自分と同じ思いをしてほしくない。」だから、フラットサポートを作りました。航空券の取り方を変えるだけで、数万円〜数十万円変わる。その事実を、もっと多くの人に知ってほしい。あなたの海外挑戦を、全力でサポートします。", 
10, False, TEXT_MEDIUM)

# ===== スライド20: よくある質問 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "FAQ", "よくある質問")

faqs = [
    ("Q. 本当に安くなりますか？", "A. はい。最大34万円お得になった実績があります。スカイスキャナーやtrip.comより安い最安値を見つけます。"),
    ("Q. 家族は何人まで使えますか？", "A. 1契約につき最大9名様まで。ご家族全員分を追加料金なしで手配します。"),
    ("Q. 解約はいつでもできますか？", "A. はい。いつでも解約可能、手数料もかかりません。初月500円で合わなければそのまま解約できます。"),
    ("Q. 旅行だけでも使えますか？", "A. もちろん。留学だけでなく、海外旅行、ワーホリ、ビジネス出張など、あらゆる海外渡航に対応。")
]

y = 1.15
for q, a in faqs:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.62, BG_LIGHT)
    add_textbox(slide, 0.55, y + 0.06, 8.9, 0.22, q, 11, True, TEXT_DARK)
    add_textbox(slide, 0.55, y + 0.32, 8.9, 0.26, a, 9, False, TEXT_MEDIUM)
    y += 0.68

# ===== スライド21: まとめ・CTA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.25, 10, 0.4, "まずは初月500円で", 26, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.65, 10, 0.4, "お試しください", 26, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.1, 10, 0.2, "航空券1回の手配で、数万円〜数十万円の差が出ます。", 12, False, WHITE, PP_ALIGN.CENTER)

price_box = add_rounded_box(slide, 3.2, 1.4, 3.6, 0.95, RGBColor(44, 82, 130))
add_textbox(slide, 3.3, 1.48, 3.4, 0.16, "初月", 9, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 3.3, 1.62, 3.4, 0.5, "500円", 44, True, ACCENT_LIGHT, PP_ALIGN.CENTER)

features = ["✓ 家族全員分OK", "✓ LINEで相談", "✓ いつでも解約可", "✓ コミュニティ"]
x = 1.0
for feat in features:
    add_textbox(slide, x, 2.5, 2.0, 0.22, feat, 10, False, WHITE, PP_ALIGN.CENTER)
    x += 2.0

add_textbox(slide, 0, 2.8, 10, 0.16, "※ 翌月以降 8,800円/月 ・ 解約手数料なし", 9, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.4, 3.1, 9.2, 0.18, "【含まれるサポート】", 10, True, WHITE)
support_list = "航空券手配 / 海外保険案内 / 24時間LINE / 海外情勢配信 / 留学サポート / ビザ代行(プレミアム) / 英会話(プレミアム) / 就職サポート / コミュニティ"
add_textbox(slide, 0.4, 3.3, 9.2, 0.3, support_list, 8, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド22: お問い合わせ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.6, 10, 0.45, "お問い合わせ", 32, True, WHITE, PP_ALIGN.CENTER)

contact_box = add_rounded_box(slide, 2.2, 1.3, 5.6, 1.9, RGBColor(44, 82, 130))

contacts = [
    ("会社名", "株式会社UNISIA（ユニシア）"),
    ("代表取締役", "井上 智羅（いのうえ とらい）"),
    ("所在地", "福岡県福岡市博多区博多駅前1丁目23番2号"),
    ("サービスURL", "ltdunisia.memberpay.jp"),
    ("公式LINE", "@unisia")
]

y = 1.45
for label, value in contacts:
    add_textbox(slide, 2.4, y, 1.8, 0.24, label, 10, False, WHITE)
    add_textbox(slide, 4.3, y, 3.3, 0.24, value, 10, False, WHITE)
    y += 0.33

add_textbox(slide, 0, 3.5, 10, 0.3, '"すべての人に、気軽な海外挑戦を"', 16, False, WHITE, PP_ALIGN.CENTER)

# 保存
output_path = "/Users/user/Library/Mobile Documents/com~apple~CloudDocs/Odsidian/トライ/02_ビジネス/ユニシア/projects/sales-material/フラットサポート営業資料_v4.pptx"
prs.save(output_path)
print(f"PowerPointファイルを作成しました: {output_path}")
print("全22スライド")
