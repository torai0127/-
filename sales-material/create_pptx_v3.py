#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
フラットサポート営業資料 PowerPoint作成スクリプト v3
- 被り削除
- サポート内容を詳しく説明（複数スライドに分割）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラー定義
PRIMARY = RGBColor(26, 54, 93)
ACCENT = RGBColor(237, 137, 54)
ACCENT_LIGHT = RGBColor(246, 173, 85)
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(26, 32, 44)
TEXT_MEDIUM = RGBColor(74, 85, 104)
TEXT_LIGHT = RGBColor(113, 128, 150)
SKY_BLUE = RGBColor(66, 153, 225)
SUCCESS = RGBColor(72, 187, 120)
DANGER = RGBColor(229, 62, 62)
BG_LIGHT = RGBColor(247, 250, 252)
BG_CREAM = RGBColor(255, 250, 240)
PINK_LIGHT = RGBColor(254, 215, 215)

def add_text_frame(shape, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans JP"
    return tf

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Noto Sans JP"
    p.alignment = align
    return box

def add_rounded_box(slide, left, top, width, height, fill_color, line_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if line_color:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box

def add_rect(slide, left, top, width, height, fill_color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    rect.line.fill.background()
    return rect

def add_section_header(slide, tag_text, title):
    tag = add_rounded_box(slide, 0.4, 0.3, len(tag_text) * 0.12 + 0.4, 0.3, SKY_BLUE)
    add_text_frame(tag, tag_text, 9, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, 0.4, 0.7, 9, 0.6, title, 28, True, PRIMARY)
    add_rect(slide, 0.4, 1.25, 9.2, 0.05, ACCENT)

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ===== スライド1: 表紙 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.3, 10, 0.4, "初月500円キャンペーン中", 14, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.8, 10, 0.3, "LTD. UNISIA", 12, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.6, 10, 0.8, "フラットサポート", 48, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 2.5, 10, 0.5, "月額制の航空券手配 & 海外挑戦サポート", 18, False, WHITE, PP_ALIGN.CENTER)

stats = [("34", "万円", "最大節約額"), ("17", "カ国", "代表の渡航経験"), ("500", "円", "初月お試し価格")]
x_positions = [1.5, 4, 6.5]
for i, (num, unit, label) in enumerate(stats):
    box = add_rounded_box(slide, x_positions[i], 3.2, 2.2, 1.3, RGBColor(44, 82, 130))
    add_textbox(slide, x_positions[i], 3.35, 2.2, 0.7, f"{num}{unit}", 32, True, ACCENT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x_positions[i], 4.0, 2.2, 0.3, label, 11, False, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0, 5.1, 10, 0.3, "株式会社UNISIA（ユニシア）", 11, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド2: こんなお悩みありませんか？ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "PROBLEM", "こんなお悩み、ありませんか？")

problems = [
    "😰 航空券、高すぎない？",
    "😩 調べるのが面倒…",
    "😱 家族全員分、いくらになる？",
    "🤔 エージェントは高いって聞く",
    "😕 本当に最安値か分からない",
    "😣 相談できる人がいない"
]

for i, problem in enumerate(problems):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 3.1
    y = 1.5 + row * 1.1
    box = add_rounded_box(slide, x, y, 2.9, 0.9, PINK_LIGHT)
    add_textbox(slide, x + 0.1, y + 0.25, 2.7, 0.5, problem, 13, True, TEXT_DARK, PP_ALIGN.CENTER)

solution = add_rounded_box(slide, 1.5, 4.0, 7, 0.65, SKY_BLUE)
add_text_frame(solution, "これらの悩み、すべて解決します。", 18, True, WHITE, PP_ALIGN.CENTER)

# ===== スライド3: フラットサポートとは（原因と解決を統合） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SERVICE", "フラットサポートとは")

# 問題提起
problem_box = add_rounded_box(slide, 0.4, 1.45, 4.4, 1.4, PINK_LIGHT)
add_textbox(slide, 0.5, 1.55, 4.2, 0.3, "💸 なぜ航空券は高い？", 14, True, DANGER, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.9, 4.2, 0.8, "多くの人はエージェントの言い値で購入。\n実は取り方次第で数万円〜数十万円変わる。", 11, False, TEXT_MEDIUM, PP_ALIGN.CENTER)

# 解決策
solution_box = add_rounded_box(slide, 5.0, 1.45, 4.6, 1.4, PRIMARY)
add_textbox(slide, 5.1, 1.55, 4.4, 0.3, "✈️ だから私たちが探します", 14, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 5.1, 1.9, 4.4, 0.8, "17カ国渡航の代表が最安値を徹底追求。\nLINEで相談するだけで全部おまかせ。", 11, False, WHITE, PP_ALIGN.CENTER)

# サービス説明
add_textbox(slide, 0.4, 3.0, 9.2, 0.35, "業界初のサブスクリプション型 航空券手配 & 海外挑戦サポートサービス", 14, True, ACCENT, PP_ALIGN.CENTER)

# フロー
flow_items = ["サブスク登録", "公式LINE追加", "LINEで相談", "サポート開始"]
x = 0.5
for i, item in enumerate(flow_items):
    box = add_rounded_box(slide, x, 3.45, 2.0, 0.45, BG_LIGHT, SKY_BLUE)
    add_text_frame(box, item, 11, True, PRIMARY, PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + 2.05, 3.55, 0.3, 0.3, "→", 16, False, ACCENT, PP_ALIGN.CENTER)
    x += 2.35

# 比較ポイント（簡潔に）
comparison = [("料金", "高額パッケージ → シンプル月額制"), ("航空券", "割高 → 最安値追求"), ("解約", "違約金あり → いつでも無料")]
y = 4.1
for item, desc in comparison:
    add_textbox(slide, 0.5, y, 9, 0.25, f"✓ {item}: {desc}", 11, False, TEXT_MEDIUM)
    y += 0.28

# ===== スライド4: 選ばれる3つの理由 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "FEATURES", "選ばれる3つの理由")

features = [
    ("💰", "POINT 01", "最大34万円安くなる", "大手エージェントより圧倒的に安く。世界一周16万円の実績。"),
    ("👨‍👩‍👧‍👦", "POINT 02", "家族全員分サポート", "月額8,800円で最大9名まで。追加料金なし。"),
    ("💬", "POINT 03", "LINEでいつでも相談", "「この時期いくら？」気軽に何度でもOK。")
]

x = 0.4
for icon, point, title, desc in features:
    box = add_rounded_box(slide, x, 1.45, 3.0, 2.0, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.45, 3.0, 0.06, ACCENT)
    add_textbox(slide, x, 1.6, 3.0, 0.4, icon, 28, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.0, 3.0, 0.2, point, 9, True, ACCENT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.25, 3.0, 0.4, title, 14, True, PRIMARY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.7, 2.8, 0.6, desc, 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    x += 3.2

bottom = add_rounded_box(slide, 0.4, 3.6, 9.2, 0.55, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 3.7, 9, 0.4, "さらに、初月500円 でお試し。合わなければ即解約OK。", 14, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド5: 実際の事例 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "RESULTS", "実際にお得になった事例")

cases = [
    ("CASE 01", "🌍 世界一周航空券", "50万円", "16万円", "34万円お得！"),
    ("CASE 02", "🇬🇧 イギリス往復", "33万円", "11万円", "22万円お得！"),
    ("CASE 03", "🇦🇺 オーストラリア往復", "18万円", "12万円", "6万円お得！")
]

x = 0.4
for case_num, dest, old_price, new_price, saving in cases:
    box = add_rounded_box(slide, x, 1.45, 3.0, 2.3, WHITE, BG_LIGHT)
    add_rect(slide, x, 1.45, 3.0, 0.06, ACCENT)
    badge = add_rounded_box(slide, x + 0.7, 1.55, 1.6, 0.28, PRIMARY)
    add_text_frame(badge, case_num, 9, True, WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x, 1.9, 3.0, 0.3, dest, 13, True, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.25, 3.0, 0.22, f"大手: {old_price}", 11, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.5, 3.0, 0.3, f"→ {new_price}", 18, True, PRIMARY, PP_ALIGN.CENTER)
    saving_box = add_rounded_box(slide, x + 0.3, 2.9, 2.4, 0.4, ACCENT)
    add_text_frame(saving_box, saving, 13, True, WHITE, PP_ALIGN.CENTER)
    x += 3.2

# 留学費用
add_textbox(slide, 0.4, 3.95, 9.2, 0.35, "📚 6ヶ月留学: 大手190万円 → フラットサポート65万円 = 約125万円お得！", 13, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0, 4.35, 10, 0.2, "※時期や条件によって価格は変動します", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# ===== スライド6: サポート内容①（航空券・手配系） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 01", "航空券・手配サポート")

supports1 = [
    ("✈️", "航空券の手配", 
     "最安値を徹底追求して手配",
     "・複数の航空会社・ルートを比較検討\n・時期による価格変動を考慮した最適なタイミングを提案\n・乗り継ぎ便も含めた総合的な比較\n・直接予約より安くなるルートを発見"),
    ("🚗", "空港送迎手配", 
     "現地到着後も安心サポート",
     "・現地空港から滞在先までの送迎を手配\n・言葉が通じなくても安心の日本語対応ドライバー\n・深夜・早朝便でも対応可能\n・留学生向けに学校までの送迎も可能"),
    ("📄", "ビザ発行代行", 
     "面倒な手続きを全て代行",
     "・必要書類の案内から申請まで全サポート\n・国ごとに異なる要件を熟知\n・申請ミスによる却下リスクを回避\n・緊急対応も可能（プレミアムプラン）")
]

y = 1.45
for icon, title, subtitle, details in supports1:
    box = add_rounded_box(slide, 0.4, y, 9.2, 1.0, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.1, 0.5, 0.4, icon, 24, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y + 0.1, 2.5, 0.3, title, 14, True, PRIMARY)
    add_textbox(slide, 1.1, y + 0.4, 2.5, 0.25, subtitle, 10, False, ACCENT)
    add_textbox(slide, 3.8, y + 0.1, 5.6, 0.8, details, 9, False, TEXT_MEDIUM)
    y += 1.1

# ===== スライド7: サポート内容②（保険・安全系） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 02", "安全・安心サポート")

supports2 = [
    ("🛡️", "海外保険案内", 
     "最適な保険プランをご提案",
     "・渡航先・期間・目的に応じた保険を比較\n・クレジットカード付帯保険との組み合わせも提案\n・保険金請求の際のサポートも対応\n・長期滞在向けの割安プランもご紹介"),
    ("🆘", "現地トラブル対応", 
     "何かあっても日本からサポート",
     "・パスポート紛失・盗難時の対応サポート\n・病院への付き添い・通訳手配\n・航空便欠航時の振替手配\n・日本からリモートで24時間対応可能"),
    ("🗣️", "現地日本語サポート", 
     "言葉の壁も安心",
     "・現地での困りごとを日本語で相談\n・病院・警察など緊急時の通訳サポート\n・各種手続きのサポート（銀行口座開設など）\n・LINEでいつでも連絡可能")
]

y = 1.45
for icon, title, subtitle, details in supports2:
    box = add_rounded_box(slide, 0.4, y, 9.2, 1.0, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.1, 0.5, 0.4, icon, 24, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y + 0.1, 2.5, 0.3, title, 14, True, PRIMARY)
    add_textbox(slide, 1.1, y + 0.4, 2.5, 0.25, subtitle, 10, False, ACCENT)
    add_textbox(slide, 3.8, y + 0.1, 5.6, 0.8, details, 9, False, TEXT_MEDIUM)
    y += 1.1

# ===== スライド8: サポート内容③（体験・アクティビティ） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 03", "体験・アクティビティ")

supports3 = [
    ("🗺️", "ツアーガイド", 
     "現地でしか味わえない体験を",
     "・会員限定の現地ツアーを定期開催\n・観光名所だけでなくローカルスポットも案内\n・他の会員との交流の機会にも\n・オプションで個別ガイドも手配可能"),
    ("📷", "海外現地撮影", 
     "思い出を形に残す",
     "・プロカメラマンによる撮影サービス\n・留学・ワーホリの記念撮影\n・SNS映えするスポットでの撮影\n・データ納品で家族にも共有可能"),
    ("💬", "24時間LINEサポート", 
     "いつでも・何度でも相談OK",
     "・時差があっても安心の24時間対応\n・些細な質問から緊急時まで対応\n・航空券以外の相談もOK（おすすめスポットなど）\n・返信は原則24時間以内")
]

y = 1.45
for icon, title, subtitle, details in supports3:
    box = add_rounded_box(slide, 0.4, y, 9.2, 1.0, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.1, 0.5, 0.4, icon, 24, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y + 0.1, 2.5, 0.3, title, 14, True, PRIMARY)
    add_textbox(slide, 1.1, y + 0.4, 2.5, 0.25, subtitle, 10, False, ACCENT)
    add_textbox(slide, 3.8, y + 0.1, 5.6, 0.8, details, 9, False, TEXT_MEDIUM)
    y += 1.1

# ===== スライド9: サポート内容④（キャリア・学習） =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "SUPPORT 04", "キャリア・学習サポート")

supports4 = [
    ("🎓", "英会話教室（月4回）", 
     "オンラインマンツーマンレッスン",
     "・プレミアムプラン限定サービス\n・ネイティブ or バイリンガル講師が担当\n・渡航前の準備から帰国後のスキル維持まで\n・日程は柔軟に調整可能"),
    ("💼", "帰国後就職サポート", 
     "海外経験を活かしたキャリア支援",
     "・求人紹介・マッチング（海外経験を活かせる企業）\n・履歴書・職務経歴書の添削\n・面接対策（海外経験のアピール方法）\n・キャリア相談（将来のキャリアプラン設計）"),
]

y = 1.45
for icon, title, subtitle, details in supports4:
    box = add_rounded_box(slide, 0.4, y, 9.2, 1.2, BG_LIGHT)
    add_textbox(slide, 0.5, y + 0.15, 0.5, 0.4, icon, 24, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, 1.1, y + 0.15, 2.5, 0.3, title, 14, True, PRIMARY)
    add_textbox(slide, 1.1, y + 0.5, 2.5, 0.25, subtitle, 10, False, ACCENT)
    add_textbox(slide, 3.8, y + 0.15, 5.6, 0.9, details, 9, False, TEXT_MEDIUM)
    y += 1.35

# サポート数まとめ
add_textbox(slide, 0.4, 4.2, 9.2, 0.35, "全10項目のサポートで、渡航前〜渡航中〜帰国後まで一貫してサポートします。", 13, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド10: コミュニティ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "COMMUNITY", "会員限定コミュニティ")

main = add_rounded_box(slide, 0.4, 1.45, 9.2, 1.0, PRIMARY)
add_textbox(slide, 0.5, 1.55, 9, 0.3, "🎮 Discordコミュニティに無料参加", 17, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.9, 9, 0.45, "フラットサポート会員だけが参加できる限定コミュニティ。\n同じ目標を持つ仲間と繋がり、情報交換や相談ができます。", 12, False, WHITE, PP_ALIGN.CENTER)

contents = [
    ("👥", "会員同士の交流", "海外経験者・これから行く人が\n集まり、リアルな情報を共有"),
    ("💡", "先輩からのアドバイス", "留学・ワーホリ経験者が\n質問に答えてくれる"),
    ("🎉", "イベント開催", "オンライン交流会・\nオフラインミートアップを定期開催")
]

x = 0.4
for icon, title, desc in contents:
    box = add_rounded_box(slide, x, 2.6, 3.0, 1.3, BG_LIGHT)
    add_textbox(slide, x, 2.7, 3.0, 0.35, icon, 24, False, TEXT_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.05, 3.0, 0.25, title, 12, True, PRIMARY, PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 3.35, 2.8, 0.5, desc, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    x += 3.2

add_textbox(slide, 0.4, 4.05, 9.2, 0.3, "一人で悩まない。仲間がいるから、海外挑戦が楽しくなる。", 13, True, ACCENT, PP_ALIGN.CENTER)

# ===== スライド11: 料金プラン =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "PRICING", "料金プラン")

# スタンダードプラン
std = add_rounded_box(slide, 0.4, 1.45, 5.5, 2.9, WHITE, ACCENT)
badge = add_rounded_box(slide, 1.8, 1.3, 2.6, 0.35, ACCENT)
add_text_frame(badge, "人気No.1 / 初月500円", 11, True, WHITE, PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 1.6, 5.3, 0.3, "スタンダードプラン", 17, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 1.9, 5.3, 0.2, "航空券手配 + フルサポート", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.1, 5.3, 0.15, "初月", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.2, 5.3, 0.5, "500円", 40, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.65, 5.3, 0.18, "翌月以降 8,800円/月", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

std_features = ["✓ 航空券の最安値検索＆手配", "✓ 家族全員分サポート（追加料金なし）", "✓ 24時間LINEサポート", "✓ Discordコミュニティ参加", "✓ 現地トラブル対応 / 海外保険案内", "✓ 帰国後就職サポート"]
y = 2.9
for feat in std_features:
    add_textbox(slide, 0.65, y, 5.1, 0.2, feat, 9, False, TEXT_MEDIUM)
    y += 0.21

# プレミアムプラン
prem = add_rounded_box(slide, 6.1, 1.45, 3.5, 2.9, WHITE, BG_LIGHT)
add_textbox(slide, 6.2, 1.6, 3.3, 0.3, "プレミアムプラン", 14, True, PRIMARY, PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 1.9, 3.3, 0.2, "留学・ワーホリをフルサポート", 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, 6.2, 2.2, 3.3, 0.4, "19,800円/月", 22, True, PRIMARY, PP_ALIGN.CENTER)

prem_features = ["✓ スタンダードの内容すべて", "✓ ビザ発行代行", "✓ 月4回オンライン英会話", "✓ 空港送迎手配", "✓ 留学フルサポート"]
y = 2.7
for feat in prem_features:
    add_textbox(slide, 6.25, y, 3.1, 0.2, feat, 9, False, TEXT_MEDIUM)
    y += 0.23

bottom = add_rounded_box(slide, 0.4, 4.45, 9.2, 0.4, BG_CREAM, ACCENT)
add_textbox(slide, 0.5, 4.5, 9, 0.3, "いつでも解約OK・解約手数料なし", 12, True, PRIMARY, PP_ALIGN.CENTER)

# ===== スライド12: 他社比較 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "COMPARE", "留学エージェントとの違い")

add_rect(slide, 0.4, 1.45, 2.6, 0.38, PRIMARY)
add_rect(slide, 3.0, 1.45, 3.2, 0.38, RGBColor(160, 174, 192))
add_rect(slide, 6.2, 1.45, 3.4, 0.38, ACCENT)
add_textbox(slide, 0.5, 1.5, 2.4, 0.28, "比較項目", 11, True, WHITE)
add_textbox(slide, 3.1, 1.5, 3.0, 0.28, "留学エージェント", 11, True, WHITE)
add_textbox(slide, 6.3, 1.5, 3.2, 0.28, "フラットサポート", 11, True, WHITE)

comparison = [
    ("対象", "留学する人のみ", "✓ 旅行でも留学でもOK"),
    ("料金体系", "パッケージ（高額）", "✓ シンプルな月額制"),
    ("航空券", "割高なことが多い", "✓ 最安値を徹底追求"),
    ("家族利用", "追加料金あり", "✓ 追加料金なし（最大9名）"),
    ("相談方法", "予約制の面談", "✓ LINEでいつでも"),
    ("解約", "違約金あり", "✓ いつでも無料"),
    ("お試し", "なし", "✓ 初月500円"),
    ("コミュニティ", "なし or 有料", "✓ 無料で参加可能")
]

y = 1.88
for item, old, new in comparison:
    add_textbox(slide, 0.5, y, 2.4, 0.26, item, 10, True, TEXT_DARK)
    add_textbox(slide, 3.1, y, 3.0, 0.26, old, 10, False, TEXT_LIGHT)
    add_textbox(slide, 6.3, y, 3.2, 0.26, new, 10, False, SUCCESS)
    y += 0.31

# ===== スライド13: お客様の声 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "VOICE", "ご利用者様の声")

voices = [
    (True, "Y.T", "20代・男性｜世界一周", "「世界一周の航空券が合計16万円で取れました！大手だと50万円近くかかると言われていたので、正直信じられなかったです。」"),
    (False, "M.S", "20代・女性｜ワーホリ", "「航空券だけで6万円浮いて、その分を現地での生活費に回せました。余計なパッケージを押し付けられないのも良かったです！」"),
    (False, "K.N", "30代・男性｜家族旅行", "「家族4人でハワイに行きました。自分で調べた価格より8万円も安く取れました。LINEで全部やってくれるのが助かりました！」"),
    (True, "S.H", "30代・女性｜イギリス留学", "「イギリス留学の航空券、大手だと33万円と言われたのが11万円で取れました。22万円も浮いて語学学校の期間を延ばせました。」")
]

y = 1.45
for featured, name, tag, text in voices:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.7, PRIMARY if featured else BG_LIGHT)
    text_color = WHITE if featured else TEXT_MEDIUM
    add_textbox(slide, 0.55, y + 0.08, 2.2, 0.25, f"{name}｜{tag}", 10, True, WHITE if featured else TEXT_DARK)
    add_textbox(slide, 0.55, y + 0.35, 8.9, 0.3, text, 10, False, text_color)
    y += 0.78

# ===== スライド14: 代表紹介 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "ABOUT", "代表紹介")

add_textbox(slide, 0.4, 1.5, 3.2, 0.3, "🌍 17カ国渡航", 14, True, ACCENT, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 1.85, 3.2, 0.4, "井上 智羅（トライ）", 18, True, TEXT_DARK, PP_ALIGN.CENTER)
add_textbox(slide, 0.4, 2.25, 3.2, 0.25, "株式会社UNISIA 代表取締役", 10, False, TEXT_LIGHT, PP_ALIGN.CENTER)

# 経歴
career = "2020年 専門学校卒業後、マルタ・NZ留学\n2021年 世界一周の旅 & ノマドフリーランス\n2021年 株式会社ファーストビュー提携（ウェブスキ設立）\n2025年 株式会社Unisia設立（フラットサポート開始）"
add_textbox(slide, 3.8, 1.5, 5.8, 1.0, career, 10, False, TEXT_MEDIUM)

# ストーリー
story_box = add_rounded_box(slide, 0.4, 2.7, 9.2, 1.5, BG_LIGHT)
add_textbox(slide, 0.5, 2.8, 9, 1.3, 
"留学エージェントに200万円損した経験があります。高額なパッケージを売りつけられ、後から「もっと安くできた」と知りました。\n\n「あの時の自分と同じ思いをしてほしくない。」だから、フラットサポートを作りました。航空券の取り方を変えるだけで、数万円〜数十万円変わる。その事実を、もっと多くの人に知ってほしい。あなたの海外挑戦を、全力でサポートします。", 
11, False, TEXT_MEDIUM)

# ===== スライド15: よくある質問 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "FAQ", "よくある質問")

faqs = [
    ("Q. 本当に安くなりますか？", "A. はい。最大34万円お得になった実績があります。常に最安値を徹底的にお探しします。"),
    ("Q. 家族は何人まで使えますか？", "A. 1契約につき最大9名様まで。ご家族全員分を追加料金なしで手配します。"),
    ("Q. 解約はいつでもできますか？", "A. はい。いつでも解約可能、手数料もかかりません。初月500円で合わなければそのまま解約できます。"),
    ("Q. 旅行だけでも使えますか？", "A. もちろん。留学だけでなく、海外旅行、ワーホリ、ビジネス出張など、あらゆる海外渡航に対応。")
]

y = 1.45
for q, a in faqs:
    box = add_rounded_box(slide, 0.4, y, 9.2, 0.7, BG_LIGHT)
    add_textbox(slide, 0.55, y + 0.08, 8.9, 0.25, q, 12, True, TEXT_DARK)
    add_textbox(slide, 0.55, y + 0.38, 8.9, 0.28, a, 10, False, TEXT_MEDIUM)
    y += 0.78

# ===== スライド16: まとめ・CTA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.5, 10, 0.6, "まずは初月500円で\nお試しください", 30, True, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.3, 10, 0.3, "航空券1回の手配で、数万円〜数十万円の差が出ます。", 14, False, WHITE, PP_ALIGN.CENTER)

price_box = add_rounded_box(slide, 3.2, 1.8, 3.6, 1.1, RGBColor(44, 82, 130))
add_textbox(slide, 3.3, 1.9, 3.4, 0.2, "初月", 11, False, WHITE, PP_ALIGN.CENTER)
add_textbox(slide, 3.3, 2.1, 3.4, 0.6, "500円", 52, True, ACCENT_LIGHT, PP_ALIGN.CENTER)

features = ["✓ 家族全員分OK", "✓ LINEでいつでも相談", "✓ いつでも解約可能", "✓ コミュニティ参加"]
x = 1.0
for feat in features:
    add_textbox(slide, x, 3.1, 2.0, 0.3, feat, 11, False, WHITE, PP_ALIGN.CENTER)
    x += 2.0

add_textbox(slide, 0, 3.55, 10, 0.2, "※ 翌月以降 8,800円/月 ・ 解約手数料なし", 10, False, WHITE, PP_ALIGN.CENTER)

# サポート一覧
add_textbox(slide, 0.4, 3.9, 9.2, 0.25, "【含まれるサポート】", 11, True, WHITE)
support_list = "航空券手配 / 空港送迎 / ビザ代行 / 海外保険案内 / 24時間LINE / 現地トラブル対応 / ツアーガイド / 英会話(プレミアム) / 就職サポート / コミュニティ"
add_textbox(slide, 0.4, 4.15, 9.2, 0.4, support_list, 9, False, WHITE, PP_ALIGN.CENTER)

# ===== スライド17: お問い合わせ =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 10, 5.625, PRIMARY)

add_textbox(slide, 0, 0.7, 10, 0.5, "お問い合わせ", 34, True, WHITE, PP_ALIGN.CENTER)

contact_box = add_rounded_box(slide, 2.2, 1.5, 5.6, 2.1, RGBColor(44, 82, 130))

contacts = [
    ("会社名", "株式会社UNISIA（ユニシア）"),
    ("代表取締役", "井上 智羅（いのうえ とらい）"),
    ("所在地", "福岡県福岡市博多区博多駅前1丁目23番2号"),
    ("サービスURL", "ltdunisia.memberpay.jp"),
    ("公式LINE", "@unisia")
]

y = 1.65
for label, value in contacts:
    add_textbox(slide, 2.4, y, 1.8, 0.26, label, 11, False, WHITE)
    add_textbox(slide, 4.3, y, 3.3, 0.26, value, 11, False, WHITE)
    y += 0.36

add_textbox(slide, 0, 3.9, 10, 0.35, '"すべての人に、気軽な海外挑戦を"', 17, False, WHITE, PP_ALIGN.CENTER)

# 保存
output_path = "/Users/user/Library/Mobile Documents/com~apple~CloudDocs/Odsidian/トライ/02_ビジネス/ユニシア/projects/sales-material/フラットサポート営業資料_v3.pptx"
prs.save(output_path)
print(f"PowerPointファイルを作成しました: {output_path}")
print("全17スライド")
