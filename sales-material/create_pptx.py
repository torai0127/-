#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
フラットサポート営業資料 PowerPoint作成スクリプト
Google スライドで開けるように .pptx 形式で出力
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# カラー定義
PRIMARY = RGBColor(26, 54, 93)  # #1a365d
ACCENT = RGBColor(237, 137, 54)  # #ed8936
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(26, 32, 44)
TEXT_MEDIUM = RGBColor(74, 85, 104)
TEXT_LIGHT = RGBColor(113, 128, 150)
SKY_BLUE = RGBColor(66, 153, 225)
SUCCESS = RGBColor(72, 187, 120)
DANGER = RGBColor(229, 62, 62)

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

def add_title_slide(prs, title, subtitle=""):
    """タイトルスライドを追加"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # タイトル
    add_text_box(slide, 0.5, 2.5, 9, 1.5, title, font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    
    if subtitle:
        add_text_box(slide, 0.5, 4, 9, 0.8, subtitle, font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    
    return slide

def add_content_slide(prs, title, content_items, section_tag=""):
    """コンテンツスライドを追加"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # セクションタグ
    if section_tag:
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(1.2), Inches(0.35))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = SKY_BLUE
        tag_box.line.fill.background()
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_tag
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # タイトル
    add_text_box(slide, 0.5, 0.9, 9, 0.8, title, font_size=28, bold=True, color=PRIMARY)
    
    # 区切り線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # コンテンツ
    y_pos = 1.9
    for item in content_items:
        add_text_box(slide, 0.5, y_pos, 9, 0.5, item, font_size=14, color=TEXT_MEDIUM)
        y_pos += 0.45
    
    return slide

# プレゼンテーション作成
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9

# ===== スライド1: 表紙 =====
slide = add_title_slide(prs, "フラットサポート", "月額制の航空券手配 & 海外挑戦サポート")
add_text_box(slide, 0.5, 0.5, 9, 0.4, "初月500円キャンペーン中", font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 1.5, 9, 0.4, "LTD. UNISIA", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# 数字セクション
stats_y = 3.2
add_text_box(slide, 1, stats_y, 2.5, 0.8, "34万円\n最大節約額", font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 4, stats_y, 2.5, 0.8, "17カ国\n代表の渡航経験", font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 7, stats_y, 2.5, 0.8, "500円\n初月お試し価格", font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 5, 9, 0.3, "株式会社UNISIA（ユニシア）", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ===== スライド2: こんなお悩みありませんか？ =====
slide = add_content_slide(prs, "こんなお悩み、ありませんか？", [], "PROBLEM")
problems = [
    "😰 航空券、高すぎない？",
    "😩 調べるのが面倒…",
    "😱 家族全員分、いくらになる？",
    "🤔 エージェントは高いって聞く",
    "😕 本当に最安値か分からない",
    "😣 相談できる人がいない"
]
y = 1.9
for i, problem in enumerate(problems):
    x = 0.5 + (i % 3) * 3.1
    if i == 3:
        y = 3.2
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(1))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(254, 215, 215)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = problem
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER

# 解決ボックス
solution_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.5), Inches(6), Inches(0.7))
solution_box.fill.solid()
solution_box.fill.fore_color.rgb = SKY_BLUE
solution_box.line.fill.background()
tf = solution_box.text_frame
p = tf.paragraphs[0]
p.text = "これらの悩み、すべて解決します。"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ===== スライド3: なぜ航空券は高いのか？ =====
slide = add_content_slide(prs, "なぜ、航空券は高いのか？", [], "REASON")

# 左側ボックス
left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(4), Inches(2.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(254, 215, 215)
left_box.line.fill.background()

add_text_box(slide, 0.7, 2.1, 3.6, 0.5, "答えはシンプル。", font_size=18, bold=True, color=DANGER, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.6, 3.6, 1, "「探し方を\n知らない」\nから。", font_size=24, bold=True, color=DANGER, align=PP_ALIGN.CENTER)

# 右側の理由リスト
reasons = [
    "💸 多くの人は大手比較サイトやエージェントの言い値で買っている",
    "🔍 実はルート・時期・航空会社で数万円〜数十万円変わる",
    "⚠️ エージェントは「最安値」ではなく利益率の高い航空券を提案"
]
y = 2.0
for reason in reasons:
    add_text_box(slide, 4.8, y, 4.7, 0.7, reason, font_size=12, color=TEXT_MEDIUM)
    y += 0.8

# 下部メッセージ
bottom_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
bottom_box.line.color.rgb = ACCENT

add_text_box(slide, 0.7, 4.7, 8.6, 0.4, "知っている人は、数十万円得している。この差を埋めるのが、フラットサポートです。", font_size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ===== スライド4: だから私たちが解決します =====
slide = add_content_slide(prs, "だから、私たちが解決します", [], "SOLUTION")

# メインボックス
main_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(9), Inches(2.2))
main_box.fill.solid()
main_box.fill.fore_color.rgb = PRIMARY
main_box.line.fill.background()

add_text_box(slide, 0.7, 2.1, 8.6, 0.6, "あなたの代わりに、最安値を探します。", font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.7, 8.6, 0.5, "17カ国渡航の経験を持つ代表が、「最安値のルート」を徹底的に探して提案します。", font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ステップ
steps = ["1. LINEで相談", "2. 最安値を探索", "3. 複数プラン提案", "4. 手配完了"]
x = 1
for step in steps:
    step_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.3), Inches(2), Inches(0.5))
    step_box.fill.solid()
    step_box.fill.fore_color.rgb = RGBColor(44, 82, 130)
    step_box.line.fill.background()
    tf = step_box.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    x += 2.2

# 下部メッセージ
add_text_box(slide, 0.5, 4.4, 9, 0.8, "「○月に○○に行きたい」とLINEで送るだけ。\nあとは全部おまかせください。", font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ===== スライド5: フラットサポートとは =====
slide = add_content_slide(prs, "フラットサポートとは", [], "SERVICE")

add_text_box(slide, 0.5, 1.9, 9, 0.4, "【業界初】サブスクリプション型の航空券手配 & 海外挑戦サポートサービス", font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.4, 9, 0.4, "エージェントのように高額なパッケージを売りつけません。", font_size=14, color=TEXT_MEDIUM, align=PP_ALIGN.CENTER)

# フロー
flow_y = 2.9
flow_items = ["サブスク登録", "→", "公式LINE追加", "→", "LINEで相談", "→", "サポート開始"]
x = 0.8
for item in flow_items:
    if item == "→":
        add_text_box(slide, x, flow_y, 0.5, 0.5, item, font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 0.6
    else:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(flow_y), Inches(1.8), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(247, 250, 252)
        box.line.color.rgb = SKY_BLUE
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER
        x += 2

# 比較表
add_text_box(slide, 0.5, 3.7, 9, 0.4, "【従来のエージェントとの違い】", font_size=14, bold=True, color=PRIMARY)

table_data = [
    ["比較項目", "従来のエージェント", "フラットサポート"],
    ["料金体系", "高額パッケージ", "✓ シンプル月額制"],
    ["航空券", "利益優先で割高", "✓ 最安値を追求"],
    ["解約", "違約金あり", "✓ いつでも無料"]
]

y = 4.1
for row in table_data:
    x = 0.5
    for i, cell in enumerate(row):
        width = 2.5 if i == 0 else 3.25
        color = TEXT_DARK if i == 0 else (TEXT_LIGHT if i == 1 else SUCCESS)
        bold = (y == 4.1)
        add_text_box(slide, x, y, width, 0.35, cell, font_size=11, bold=bold, color=color)
        x += width
    y += 0.35

# ===== スライド6: 選ばれる3つの理由 =====
slide = add_content_slide(prs, "選ばれる3つの理由", [], "FEATURES")

features = [
    ("💰", "POINT 01", "最大34万円\n安くなる", "大手エージェントより圧倒的に安く航空券を手配。世界一周16万円の実績あり。"),
    ("👨‍👩‍👧‍👦", "POINT 02", "家族全員分\nサポート", "月額8,800円で、ご家族全員の航空券を手配。追加料金なし。最大9名まで。"),
    ("💬", "POINT 03", "LINEで\nいつでも相談", "「この時期はいくら？」「どのルートが安い？」気軽にLINEで何度でも相談OK。")
]

x = 0.5
for icon, point, title, desc in features:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(2.9), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(226, 232, 240)
    
    add_text_box(slide, x + 0.2, 2.0, 2.5, 0.5, icon, font_size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.5, 2.5, 0.3, point, font_size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.8, 2.5, 0.6, title, font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.4, 2.5, 0.8, desc, font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    
    x += 3.1

# 下部メッセージ
bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
bottom.fill.solid()
bottom.fill.fore_color.rgb = RGBColor(255, 250, 240)
bottom.line.color.rgb = ACCENT
add_text_box(slide, 0.7, 4.65, 8.6, 0.5, "さらに、初月500円 でお試しいただけます。合わなければ即解約OK。", font_size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ===== スライド7: 実際の事例 =====
slide = add_content_slide(prs, "実際にお得になった事例", [], "RESULTS")

cases = [
    ("CASE 01", "🌍 世界一周航空券", "50万円", "16万円", "34万円お得！"),
    ("CASE 02", "🇬🇧 イギリス往復", "33万円", "11万円", "22万円お得！"),
    ("CASE 03", "🇦🇺 オーストラリア往復", "18万円", "12万円", "6万円お得！")
]

x = 0.5
for case_num, dest, old_price, new_price, saving in cases:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(2.9), Inches(2.8))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(226, 232, 240)
    
    # 上部のオレンジライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.9), Inches(2.9), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    add_text_box(slide, x + 0.2, 2.1, 2.5, 0.3, case_num, font_size=10, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.4, 2.5, 0.4, dest, font_size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 2.85, 2.5, 0.3, f"大手: {old_price}", font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.15, 2.5, 0.4, f"→ {new_price}", font_size=20, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    
    # 節約額バッジ
    saving_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.3), Inches(3.6), Inches(2.3), Inches(0.5))
    saving_box.fill.solid()
    saving_box.fill.fore_color.rgb = ACCENT
    saving_box.line.fill.background()
    tf = saving_box.text_frame
    p = tf.paragraphs[0]
    p.text = saving
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    x += 3.1

add_text_box(slide, 0.5, 4.8, 9, 0.3, "※時期や条件によって価格は変動します", font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# ===== スライド8: 留学費用の比較 =====
slide = add_content_slide(prs, "留学費用も大幅削減", [], "RESULTS")

# 左側（大手）
left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(4), Inches(2))
left.fill.solid()
left.fill.fore_color.rgb = RGBColor(254, 215, 215)
left.line.fill.background()
add_text_box(slide, 0.7, 2.2, 3.6, 0.3, "大手留学エージェント", font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.6, 3.6, 0.6, "190万円〜", font_size=36, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 3.3, 3.6, 0.3, "6ヶ月留学の場合", font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# 右側（フラットサポート）
right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(2), Inches(4.5), Inches(2))
right.fill.solid()
right.fill.fore_color.rgb = PRIMARY
right.line.fill.background()
add_text_box(slide, 5.2, 2.2, 4.1, 0.3, "フラットサポート", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 5.2, 2.6, 4.1, 0.6, "65万円〜", font_size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 5.2, 3.3, 4.1, 0.3, "6ヶ月留学の場合", font_size=10, color=WHITE, align=PP_ALIGN.CENTER)

# 節約額
saving_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.3), Inches(4), Inches(0.7))
saving_box.fill.solid()
saving_box.fill.fore_color.rgb = ACCENT
saving_box.line.fill.background()
tf = saving_box.text_frame
p = tf.paragraphs[0]
p.text = "約125万円お得！"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ===== スライド9: サポート内容 =====
slide = add_content_slide(prs, "充実のサポート内容", [], "SUPPORT")

add_text_box(slide, 0.5, 1.9, 9, 0.3, "航空券手配だけでなく、海外挑戦に必要な10項目のサポートを提供しています。", font_size=13, color=TEXT_MEDIUM)

supports = [
    ("✈️ 航空券の手配", "最安値を徹底追求"),
    ("🚗 空港送迎手配", "到着後も安心"),
    ("📄 ビザ発行代行", "面倒な手続きを代行"),
    ("🛡️ 海外保険案内", "最適なプランをご提案"),
    ("💬 24時間LINEサポート", "いつでも何度でも相談OK"),
    ("🆘 現地トラブル対応", "渡航中の問題も対応"),
    ("🗣️ 現地日本語サポート", "言葉の壁も安心"),
    ("🗺️ ツアーガイド", "現地でのツアーを開催"),
    ("🎓 英会話教室（月4回）", "プレミアムプラン限定"),
    ("💼 帰国後就職サポート", "キャリアもサポート")
]

y = 2.3
for i, (title, desc) in enumerate(supports):
    x = 0.5 if i % 2 == 0 else 5
    if i > 0 and i % 2 == 0:
        y += 0.45
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.4), Inches(0.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 250, 252)
    box.line.fill.background()
    add_text_box(slide, x + 0.1, y + 0.05, 4.2, 0.3, f"{title} - {desc}", font_size=11, color=TEXT_MEDIUM)

# ===== スライド10: 料金プラン =====
slide = add_content_slide(prs, "料金プラン", [], "PRICING")

# スタンダードプラン
std_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(5.5), Inches(3.2))
std_box.fill.solid()
std_box.fill.fore_color.rgb = WHITE
std_box.line.color.rgb = ACCENT

# 人気No.1バッジ
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.75), Inches(2.5), Inches(0.35))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT
badge.line.fill.background()
tf = badge.text_frame
p = tf.paragraphs[0]
p.text = "人気No.1 / 初月500円"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

add_text_box(slide, 0.7, 2.2, 5.1, 0.4, "スタンダードプラン", font_size=20, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.6, 5.1, 0.3, "航空券手配 + フルサポート", font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 2.95, 5.1, 0.3, "初月", font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 3.2, 5.1, 0.6, "500円", font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.7, 3.8, 5.1, 0.3, "翌月以降 8,800円/月", font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

std_features = ["✓ 航空券の最安値検索＆手配", "✓ 家族全員分サポート（追加料金なし）", "✓ 24時間LINEサポート", "✓ 現地トラブル対応", "✓ 海外保険案内"]
y = 4.15
for feat in std_features:
    add_text_box(slide, 0.9, y, 5, 0.25, feat, font_size=10, color=TEXT_MEDIUM)
    y += 0.25

# プレミアムプラン
prem_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(1.9), Inches(3.3), Inches(3.2))
prem_box.fill.solid()
prem_box.fill.fore_color.rgb = WHITE
prem_box.line.color.rgb = RGBColor(226, 232, 240)

add_text_box(slide, 6.4, 2.2, 2.9, 0.4, "プレミアムプラン", font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
add_text_box(slide, 6.4, 2.55, 2.9, 0.3, "留学・ワーホリをフルサポート", font_size=9, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, 6.4, 2.9, 2.9, 0.5, "19,800円/月", font_size=28, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

prem_features = ["✓ スタンダードの内容すべて", "✓ ビザ発行代行", "✓ 月4回英会話教室", "✓ 留学フルサポート"]
y = 3.5
for feat in prem_features:
    add_text_box(slide, 6.4, y, 2.9, 0.25, feat, font_size=10, color=TEXT_MEDIUM)
    y += 0.3

# ===== スライド11: 他社比較 =====
slide = add_content_slide(prs, "留学エージェントとの違い", [], "COMPARE")

comparison = [
    ["比較項目", "留学エージェント", "フラットサポート"],
    ["対象", "留学する人のみ", "✓ 旅行でも留学でもOK"],
    ["料金体系", "パッケージ（高額）", "✓ シンプルな月額制"],
    ["航空券", "割高なことが多い", "✓ 最安値を徹底追求"],
    ["家族利用", "追加料金あり", "✓ 追加料金なし"],
    ["相談方法", "予約制の面談", "✓ LINEでいつでも"],
    ["解約", "違約金あり", "✓ いつでも無料"],
    ["お試し", "なし", "✓ 初月500円"]
]

y = 1.9
for row in comparison:
    x = 0.5
    for i, cell in enumerate(row):
        width = 2.5 if i == 0 else 3.25
        bg_color = PRIMARY if y == 1.9 else (RGBColor(247, 250, 252) if i == 0 else WHITE)
        text_color = WHITE if y == 1.9 else (TEXT_DARK if i == 0 else (TEXT_LIGHT if i == 1 else SUCCESS))
        
        if y == 1.9:
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.4))
            box.fill.solid()
            box.fill.fore_color.rgb = PRIMARY if i != 2 else ACCENT
            box.line.fill.background()
        
        add_text_box(slide, x + 0.1, y + 0.05, width - 0.2, 0.3, cell, font_size=11, bold=(y == 1.9), color=text_color)
        x += width
    y += 0.38

# ===== スライド12: お客様の声1 =====
slide = add_content_slide(prs, "ご利用者様の声", [], "VOICE")

voices1 = [
    ("Y", "male", "Y.Tさん（20代・男性）", "世界一周で利用", "「世界一周の航空券が合計16万円で取れました！大手だと50万円近くかかると言われていたので、正直信じられなかったです。」", True),
    ("M", "female", "M.Sさん（20代・女性）", "ワーホリ準備で利用", "「航空券だけで6万円浮いて、その分を現地での生活費に回せました。余計なパッケージを押し付けられないのも良かったです！」", False)
]

y = 1.9
for initial, gender, name, tag, text, featured in voices1:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY if featured else RGBColor(247, 250, 252)
    box.line.fill.background()
    
    text_color = WHITE if featured else TEXT_MEDIUM
    name_color = WHITE if featured else TEXT_DARK
    
    add_text_box(slide, 0.7, y + 0.1, 2, 0.3, f"{initial} | {name}", font_size=12, bold=True, color=name_color)
    add_text_box(slide, 0.7, y + 0.4, 2, 0.2, tag, font_size=10, color=ACCENT if featured else TEXT_LIGHT)
    add_text_box(slide, 0.7, y + 0.65, 8.5, 0.6, text, font_size=11, color=text_color)
    
    y += 1.4

# ===== スライド13: お客様の声2 =====
slide = add_content_slide(prs, "ご利用者様の声", [], "VOICE")

voices2 = [
    ("K", "male", "K.Nさん（30代・男性）", "家族旅行で利用", "「家族4人でハワイに行きました。自分で調べた価格より8万円も安く取れました。子どもがいると調べる時間もないので、LINEで相談するだけで全部やってくれるのが助かりました！」", False),
    ("S", "female", "S.Hさん（30代・女性）", "イギリス留学で利用", "「イギリス留学の航空券、大手だと33万円と言われたのが11万円で取れました。22万円も浮いたおかげで、現地での語学学校の期間を延ばせました。」", False)
]

y = 1.9
for initial, gender, name, tag, text, featured in voices2:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(1.3))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY if featured else RGBColor(247, 250, 252)
    box.line.fill.background()
    
    text_color = WHITE if featured else TEXT_MEDIUM
    name_color = WHITE if featured else TEXT_DARK
    
    add_text_box(slide, 0.7, y + 0.1, 2, 0.3, f"{initial} | {name}", font_size=12, bold=True, color=name_color)
    add_text_box(slide, 0.7, y + 0.4, 2, 0.2, tag, font_size=10, color=ACCENT if featured else TEXT_LIGHT)
    add_text_box(slide, 0.7, y + 0.65, 8.5, 0.6, text, font_size=11, color=text_color)
    
    y += 1.4

# メッセージ
msg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
msg_box.fill.solid()
msg_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
msg_box.line.color.rgb = ACCENT
add_text_box(slide, 0.7, 4.7, 8.6, 0.4, "多くの方に選ばれています。次はあなたが、お得に海外へ行く番です。", font_size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

# ===== スライド14: 代表紹介 =====
slide = add_content_slide(prs, "代表紹介", [], "ABOUT")

add_text_box(slide, 0.5, 1.9, 3, 0.4, "🌍 17カ国渡航", font_size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.4, 3, 0.5, "井上 智羅（トライ）", font_size=22, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.9, 3, 0.3, "株式会社UNISIA 代表取締役", font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

story = """2020年、専門学校卒業後にマルタ・NZ留学へ。
その後、世界一周の旅に出てノマドフリーランスとして活動。

一方で、留学エージェントに200万円損した経験も。
高額なパッケージを売りつけられ、後から
「もっと安くできた」と知りました。

「あの時の自分と同じ思いをしてほしくない。」
だから、フラットサポートを作りました。"""

add_text_box(slide, 3.5, 1.9, 6, 2.8, story, font_size=12, color=TEXT_MEDIUM)

# ===== スライド15: よくある質問1 =====
slide = add_content_slide(prs, "よくある質問", [], "FAQ")

faqs1 = [
    ("Q. 本当に安くなりますか？", "A. はい。大手エージェントと比較して、最大34万円お得になった実績があります。常に最安値を徹底的にお探しします。"),
    ("Q. 家族は何人まで使えますか？", "A. 1契約につき最大9名様までご利用いただけます。ご家族全員分を追加料金なしで手配します。"),
    ("Q. 解約はいつでもできますか？", "A. はい。いつでも解約可能です。解約手数料もかかりません。")
]

y = 1.9
for q, a in faqs1:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 250, 252)
    box.line.fill.background()
    
    add_text_box(slide, 0.7, y + 0.1, 8.6, 0.3, q, font_size=14, bold=True, color=TEXT_DARK)
    add_text_box(slide, 0.7, y + 0.45, 8.6, 0.4, a, font_size=11, color=TEXT_MEDIUM)
    y += 1.0

# ===== スライド16: よくある質問2 =====
slide = add_content_slide(prs, "よくある質問", [], "FAQ")

faqs2 = [
    ("Q. 航空券以外も相談できますか？", "A. はい。ビザ申請や海外保険のご案内、現地トラブル対応、治安情報の配信など、LINEでいつでもご相談いただけます。"),
    ("Q. 旅行だけでも使えますか？", "A. もちろんです。留学だけでなく、海外旅行、ワーホリ、ビジネス出張など、あらゆる海外渡航に対応しています。"),
    ("Q. 初月500円の後、自動で課金されますか？", "A. はい、翌月から8,800円/月となります。ただし、いつでも解約可能で解約手数料もありません。")
]

y = 1.9
for q, a in faqs2:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(9), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 250, 252)
    box.line.fill.background()
    
    add_text_box(slide, 0.7, y + 0.1, 8.6, 0.3, q, font_size=14, bold=True, color=TEXT_DARK)
    add_text_box(slide, 0.7, y + 0.45, 8.6, 0.4, a, font_size=11, color=TEXT_MEDIUM)
    y += 1.0

# ===== スライド17: まとめ・CTA =====
slide = add_title_slide(prs, "まずは初月500円で\nお試しください", "航空券1回の手配で、数万円〜数十万円の差が出ます。")

# 価格ボックス
price_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2.8), Inches(4), Inches(1.2))
price_box.fill.solid()
price_box.fill.fore_color.rgb = RGBColor(44, 82, 130)
price_box.line.fill.background()

add_text_box(slide, 3.2, 2.9, 3.6, 0.3, "初月", font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 3.2, 3.2, 3.6, 0.7, "500円", font_size=56, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# 特徴
features_text = "✓ 家族全員分OK　　✓ LINEでいつでも相談　　✓ いつでも解約可能"
add_text_box(slide, 0.5, 4.2, 9, 0.4, features_text, font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 4.7, 9, 0.3, "※ 翌月以降 8,800円/月 ・ 解約手数料なし", font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ===== スライド18: お問い合わせ =====
slide = add_title_slide(prs, "お問い合わせ", "")

# 連絡先ボックス
contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(6), Inches(2.2))
contact_box.fill.solid()
contact_box.fill.fore_color.rgb = RGBColor(44, 82, 130)
contact_box.line.fill.background()

contacts = [
    ("会社名", "株式会社UNISIA（ユニシア）"),
    ("代表取締役", "井上 智羅（いのうえ とらい）"),
    ("所在地", "福岡県福岡市博多区博多駅前1丁目23番2号"),
    ("サービスURL", "ltdunisia.memberpay.jp"),
    ("公式LINE", "@unisia")
]

y = 2.15
for label, value in contacts:
    add_text_box(slide, 2.3, y, 1.8, 0.3, label, font_size=11, color=WHITE)
    add_text_box(slide, 4.2, y, 3.5, 0.3, value, font_size=11, color=WHITE)
    y += 0.38

add_text_box(slide, 0.5, 4.5, 9, 0.4, '"すべての人に、気軽な海外挑戦を"', font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

# 保存
output_path = "/Users/user/Library/Mobile Documents/com~apple~CloudDocs/Odsidian/トライ/02_ビジネス/ユニシア/projects/sales-material/フラットサポート営業資料.pptx"
prs.save(output_path)
print(f"PowerPointファイルを作成しました: {output_path}")
